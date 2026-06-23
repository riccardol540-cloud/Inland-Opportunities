"""
INLAND Tunisia — Presentation Builder
Generates the .pptx pitch for the Institut français Mediterranean photography
support (publication + production strands). Mirrors the INLAND house style.
Run from this directory: `python3 build_presentation.py`
"""

from dataclasses import dataclass


@dataclass
class Story:
    headline: str
    subtitle: str
    hook: str
    proposal: str
    category: str
    tags: list
    pull_quote: str
    places: str
    reference: str


INTRO = {
    "about": (
        "INLAND is a yearly publication documenting the soul of countries "
        "before they change. Each issue embeds with named, locatable custodians "
        "of living traditions — artisans, pastoralists, cooks, singers, border "
        "families — whose stories have not yet been narrated internationally. "
        "Previous issues have covered the Faroe Islands, Georgia and Italy’s "
        "minor islands."
    ),
    "about_method": (
        "The editorial method is slow, place-based and built on trust. We spend "
        "months in pre-production research, identify specific custodians and "
        "communities, then commission photographers, writers and filmmakers to "
        "embed with them over extended field periods. The result is long-form "
        "visual storytelling grounded in named people and real places."
    ),
    "research": (
        "Tunisia is internationally legible through three loud stories — the "
        "2011 revolution that began the Arab Spring, the Mediterranean migration "
        "crossings that leave from Sfax and Zarzis, and the beach-resort coast. "
        "All three are loud; all three crowd out the quieter country. INLAND asks "
        "one question across every lead: what survived? What persists on the "
        "shore and the islands that has been absorbed by neither the "
        "Phoenician–Roman antiquity the tour buses visit, the French "
        "protectorate, the resort economy, nor the post-revolution turbulence — "
        "and can still be witnessed and photographed today."
    ),
    "research_scope": (
        "Our desk research identified 17 story leads across craft, ritual, "
        "fishing, food, language and landscape. From these, twelve storylines "
        "have been selected for editorial development — each centred on a "
        "practice that can still be witnessed, a specific place, and a named "
        "custodian or custodian community on Tunisia’s Mediterranean shores."
    ),
    "angle": (
        "We read Tunisia as a Mediterranean palimpsest — Phoenician, Roman, "
        "Amazigh, Andalusian, Jewish, Ibadi, sub-Saharan and Ottoman — through "
        "the people who keep its endangered knowledge: the palm-frond charfia "
        "fences of Kerkennah, the hand-built pottery of the women of Sejnane, the "
        "Stambeli trance liturgy of Black Tunisia. The sea is crossing, frontier, "
        "larder and memory. The revolution and the migration crisis are context, "
        "never the lens; the state is never cast as villain."
    ),
    "angle_why": (
        "This is the Mediterranean photographic work the Institut français "
        "support exists to fund: a finished INLAND edition on the monde "
        "méditerranéen, in the register of Issue 03 (Isole Minori). The "
        "publication strand covers the edition’s editing and printing; the "
        "production strand funds the field shoot across three coastal and island "
        "clusters — made with the custodians, not about them."
    ),
    "categories": [
        ("Arts & Culture", 3),
        ("Food", 2),
        ("Local communities", 2),
        ("Myth & Religion", 2),
        ("Geopolitics", 1),
        ("History", 1),
        ("Nature", 1),
    ],
}

STORIES = [
    Story(
        headline="The Palm in the Sea",
        subtitle="The Charfia of the Kerkennah Islands",
        hook=(
            "On the low, palm-fringed Kerkennah Islands off Sfax, fishermen still "
            "plant date-palm fronds in rows across the shallows to build the "
            "charfia — a great V of woven palm that channels fish on the ebb tide "
            "into capture chambers, the grounds shared by a draw of lots and "
            "worked to a maritime calendar. UNESCO inscribed it in 2020."
        ),
        proposal=(
            "INLAND follows one charfia family through the autumn setting of the "
            "fronds and a dawn lift of the traps — a fishing architecture built "
            "entirely from the palm, standing in the Mediterranean. The edition’s "
            "signature recurring image."
        ),
        category="Nature",
        tags=["UNESCO", "Endangered"],
        pull_quote="A fence of palm leaves, planted in the sea, older than any net.",
        places="Kerkennah Islands · Sfax",
        reference="UNESCO ICH 2020 (Charfia fishing); Earth Journalism Network; New Arab",
    ),
    Story(
        headline="The Women Who Build with Earth",
        subtitle="Sejnane’s Amazigh Potters",
        hook=(
            "In the Mogods hills behind Bizerte, the women of Sejnane hand-build "
            "pots, animals and the famous Sejnane dolls from wadi clay with no "
            "wheel and no kiln, firing them in open ground and painting them with "
            "two-tone geometric signs from tattoo and Berber weaving — a craft "
            "passed mother to daughter for some three thousand years."
        ),
        proposal=(
            "INLAND portraits one named potter through a full cycle — digging and "
            "kneading the clay, shaping by hand, the open firing, the painting — a "
            "Mediterranean ceramics tradition that predates the wheel and belongs "
            "entirely to women. UNESCO-inscribed 2018."
        ),
        category="Arts & Culture",
        tags=["UNESCO", "Women"],
        pull_quote="No wheel, no kiln — only the hand, the fire, and three thousand years.",
        places="Sejnane · Mogods · Bizerte",
        reference="UNESCO ICH 2018 (Pottery of the women of Sejnane); The Arab Weekly; phys.org",
    ),
    Story(
        headline="The Island of Three Prayers",
        subtitle="Djerba’s Ibadis & El Ghriba",
        hook=(
            "On Djerba, Ibadi Muslims built low, fortified, sometimes "
            "half-underground mosques across an island of menzel farmsteads and "
            "rain-cisterns, while at El Ghriba — held to be Africa’s oldest "
            "synagogue — Jewish pilgrims still gather each spring for the Lag "
            "BaOmer procession. UNESCO listed the whole island in 2023."
        ),
        proposal=(
            "INLAND works the island as one fabric — a fortified Ibadi mosque at "
            "prayer, a menzel courtyard, the El Ghriba pilgrimage — documenting "
            "the only place in the Arab world where Ibadi, Maliki, Hanafi, "
            "Christian and Jewish communities have shared a small island for "
            "centuries."
        ),
        category="Myth & Religion",
        tags=["UNESCO", "Coexistence"],
        pull_quote="Three kinds of mosque, a church and a synagogue, on one waterless island.",
        places="Djerba · Hara Kebira · Erriadh",
        reference="UNESCO WHC 2023 (Djerba); El Ghriba (Wikipedia); The Arab Weekly",
    ),
    Story(
        headline="The Villages That Speak the Old Tongue",
        subtitle="The Amazigh Ksour of the Dahar",
        hook=(
            "On the ridges of the Dahar above Tataouine, the Amazigh villages of "
            "Chenini and Douiret cling to the rock around honeycombed ghorfa "
            "granaries and a white mosque, and a dwindling number of elders still "
            "speak Tamazight where a generation ago far more did. Below, families "
            "have moved to ‘new’ villages with roads and water."
        ),
        proposal=(
            "INLAND follows a family that still keeps a hill house and a ghorfa — "
            "the olive-and-barley year, the granary, the Tamazight spoken at home "
            "— documenting one of the last living Berber-speaking pockets of a "
            "country that mostly forgot it was Amazigh."
        ),
        category="Local communities",
        tags=["Amazigh", "Language"],
        pull_quote="A language kept like seed-grain, in a granary on a hill.",
        places="Chenini · Douiret · Tataouine",
        reference="Carthage Magazine (Amazigh heritage); WildyNess; Discover Tunisia",
    ),
    Story(
        headline="The Clock That Runs Backward",
        subtitle="Testour’s Andalusians & the Malouf",
        hook=(
            "Testour, in the Medjerda valley, was rebuilt in the 17th century by "
            "Moriscos expelled from Spain; its Great Mosque carries a clock whose "
            "hands turn counter-clockwise — read locally as the exiles’ longing to "
            "wind time back toward Andalusia — and the town keeps the malouf, the "
            "Arab-Andalusian music carried across the sea."
        ),
        proposal=(
            "INLAND spends the malouf festival with a family of Testour musicians "
            "and follows the Andalusian thread through the town — the backward "
            "clock, the brick mosque, the pomegranate orchards — a four-century-old "
            "Mediterranean exile still audible."
        ),
        category="Arts & Culture",
        tags=["Andalusian", "Music"],
        pull_quote="They set the clock to run home, across a sea they could not recross.",
        places="Testour · Medjerda · Béja",
        reference="The Arab Weekly; The National; Tunisian malouf (Wikipedia)",
    ),
    Story(
        headline="The Oasis by the Sea",
        subtitle="Gabès, the Last Maritime Oasis",
        hook=(
            "Gabès holds the Mediterranean’s only seaside oasis — a three-tier "
            "garden of date palm, pomegranate, henna and vegetable running down to "
            "the shore, watered for millennia — now poisoned at its edge by the "
            "state phosphate works that have piled phosphogypsum into the gulf for "
            "fifty years. Farmers describe henna and pomegranate failing."
        ),
        proposal=(
            "INLAND documents the oasis from inside — a henna grower, the layered "
            "garden, the old irrigation — against the chimneys at its margin: a "
            "living agricultural landscape and the industry pressing on it. "
            "Tension, not indictment; the gardeners are the subject."
        ),
        category="Geopolitics",
        tags=["Ecology", "Henna"],
        pull_quote="A garden that has met the sea for three thousand years, now meeting the factory.",
        places="Gabès · Chenini Gabès · Gulf of Gabès",
        reference="Med-O-Med (Gabès oasis); Coda Story; TIMEP; SEVENSEAS",
    ),
    Story(
        headline="The Spirits Carried North",
        subtitle="Stambeli, Black Tunisia’s Healing Trance",
        hook=(
            "In a few houses in Tunis — above all Dar Barnou, once a refuge for "
            "freed slaves — the descendants of sub-Saharan people brought across "
            "the Sahara keep Stambeli, a night music of the gambri lute and iron "
            "shqasheq that calls a pantheon of sub-Saharan spirits and Muslim "
            "saints to heal through trance. It survives in a handful of "
            "brotherhoods."
        ),
        proposal=(
            "INLAND sits a full Stambeli night with one of the last troupes — the "
            "tuning, the incense, the colours each spirit demands, the dance — "
            "documenting Black Tunisia’s living liturgy with consent, as ritual "
            "rather than spectacle."
        ),
        category="Myth & Religion",
        tags=["Ritual", "Minority"],
        pull_quote="The drum crossed the desert in chains, and it is still healing.",
        places="Tunis · Dar Barnou · Sidi Saad",
        reference="Stambeli (Wikipedia); Jankowsky, Stambeli (U. Chicago); Pan-African Music",
    ),
    Story(
        headline="The Last Matanza",
        subtitle="The Bluefin Trap of Cap Bon",
        hook=(
            "At Sidi Daoud, on the tip of Cap Bon, fishermen set the madrague — a "
            "labyrinth of nets the locals call the matanza — to trap migrating "
            "Atlantic bluefin between May and July, a Phoenician-old technique that "
            "ends in a churning chamber of tuna hauled by hand. Quotas, industrial "
            "fleets and warming water have all but stilled it."
        ),
        proposal=(
            "INLAND embeds for a season at one of the Mediterranean’s last working "
            "tuna traps — the building of the net-rooms, the watch for the school, "
            "the matanza itself — a violent, ancient sea-harvest at the edge of "
            "disappearing."
        ),
        category="Food",
        tags=["Sea", "Endangered"],
        pull_quote="A maze of nets the Phoenicians would recognise, sprung once a year.",
        places="Sidi Daoud · Cap Bon · El Haouaria",
        reference="INSTM Bulletin (madrague Sidi Daoud); Mattanza (Wikipedia); harissa.com",
    ),
    Story(
        headline="White Gold",
        subtitle="Distilling Neroli at Nabeul",
        hook=(
            "Each April and May around Nabeul, families pick the bitter-orange "
            "blossom by hand before dawn, when the scent is strongest, and distil "
            "it over slow fire in copper alembics into zhar — orange-blossom water "
            "for sweets and grief and celebration — and the precious neroli oil. A "
            "million flowers make a kilogram of oil."
        ),
        proposal=(
            "INLAND follows one family’s blossom week — the pre-dawn picking, the "
            "qattara still, the bottling — a Mediterranean perfume-and-kitchen "
            "tradition where the whole coast smells of one tree for a fortnight a "
            "year."
        ),
        category="Food",
        tags=["Distillation", "Season"],
        pull_quote="For two weeks the whole coast smells of one flower.",
        places="Nabeul · Hammamet · Cap Bon",
        reference="Africanews (Nabeul orange blossom); Middle East Eye; Biolandes",
    ),
    Story(
        headline="The Gold-Threaded Loom",
        subtitle="Wedding Silk of Mahdia",
        hook=(
            "In the sea-walled medina of Mahdia, weavers still work tall "
            "draw-looms — a technology that came across the Mediterranean from "
            "Iberia — to make narrow silk strips threaded with gold and silver for "
            "the keswa el-kbira, the layered wedding dress that tells a Sahel "
            "woman’s married status. The craft was brought by Jewish artisans."
        ),
        proposal=(
            "INLAND portraits a master weaver and a bride’s family through the "
            "making and wearing of a gold-silk keswa — the loom, the shuttle, the "
            "dressing — a coastal textile tradition where a wedding is still woven "
            "by hand."
        ),
        category="Arts & Culture",
        tags=["Weaving", "Women"],
        pull_quote="A marriage you can read in the gold of the cloth.",
        places="Mahdia · the Sahel coast",
        reference="Carthage Magazine (Tunisian wedding); Discover Tunisia; shirahime.ch",
    ),
    Story(
        headline="The Knot and the Sweet",
        subtitle="The Holy City’s Carpets",
        hook=(
            "Kairouan, Islam’s fourth holy city and its first in Africa, is also "
            "Tunisia’s carpet capital, where women knot the zerbia — hundreds of "
            "thousands of knots a square metre — in workshops around the Great "
            "Mosque, and bake the date-and-semolina makroudh the city claims as "
            "its own. Both loom and oven are pressed by cheaper imports."
        ),
        proposal=(
            "INLAND works the medina with one weaving family and one makroudh "
            "house — the dyed wool, the knotting frame, the honeyed pastry — "
            "documenting the holy city through the women’s hands that keep its two "
            "famous crafts."
        ),
        category="History",
        tags=["Carpets", "Holy city"],
        pull_quote="Half a million knots, said like a prayer, in the city of the first mosque.",
        places="Kairouan · the Medina",
        reference="Discover Tunisia (weaving & carpets); Folk Culture (Kairouan carpet); WildyNess",
    ),
    Story(
        headline="The Green Mountain",
        subtitle="Cork Oak and the Khroumir",
        hook=(
            "In the Kroumirie range above Tabarka, where it snows in winter and "
            "the cork-oak forest is unlike anywhere else in Tunisia, the Khroumir "
            "mountain people harvest cork on its long bark-cycle, herd and gather "
            "among the trees, and keep a green, wet, half-forgotten north that "
            "contradicts the country’s Saharan image."
        ),
        proposal=(
            "INLAND follows a cork-harvesting family through a summer stripping "
            "season and the forest year — the orange-stripped trunks, the herding, "
            "the highland village — documenting the Mediterranean Tunisia of oak "
            "and rain rather than palm and dune."
        ),
        category="Local communities",
        tags=["Forest", "Cork"],
        pull_quote="A Tunisia of oak bark and snow, where you expected only sand.",
        places="Aïn Draham · Kroumirie · Tabarka",
        reference="Carthage Magazine (Aïn Draham); Tunisi.info (Kroumirie); Take Your Backpack",
    ),
]

TRACKS = [
    {
        "name": "The South & the Islands",
        "color": (0x1B, 0x49, 0x65),
        "stories": "Stories 01 + 03 + 04 + 06",
        "subtitle": "Sea-fence, island, ksar & oasis",
        "desc": (
            "The Mediterranean south: the charfia palm-fences of Kerkennah, the "
            "coexistence-island of Djerba (Ibadi mosques + El Ghriba), the Amazigh "
            "ksour of the Dahar, and the threatened seaside oasis of Gabès. One "
            "southern deployment."
        ),
        "geography": "Kerkennah/Sfax → Gabès → Djerba → Tataouine",
        "season": "Charfia autumn-set; El Ghriba pilgrimage (spring)",
        "story_count": "4 stories, 1 southern base",
    },
    {
        "name": "Cap Bon & the Sahel",
        "color": (0x3A, 0x5A, 0x40),
        "stories": "Stories 08 + 09 + 10 + 11",
        "subtitle": "Tuna, blossom, silk & the knot",
        "desc": (
            "The eastern shoulder: the Sidi Daoud bluefin matanza and the Nabeul "
            "orange-blossom stills of Cap Bon, the wedding silk of Mahdia and the "
            "carpets and makroudh of holy Kairouan. Sea and craft on the Sahel."
        ),
        "geography": "Cap Bon (Sidi Daoud/Nabeul) → Kairouan → Mahdia",
        "season": "Blossom Apr–May; matanza May–Jul",
        "story_count": "4 stories, spring window",
    },
    {
        "name": "The Green North",
        "color": (0x5B, 0x2C, 0x6F),
        "stories": "Stories 02 + 05 + 07 + 12",
        "subtitle": "Clay, malouf, trance & cork",
        "desc": (
            "The wet, layered north: the women potters of Sejnane, the Andalusian "
            "malouf of Testour, the Stambeli brotherhoods of Tunis, and the "
            "Kroumirie cork-oak country above Tabarka. The Tunisia of oak, exile "
            "and rain."
        ),
        "geography": "Tunis → Testour (Medjerda) → Sejnane (Bizerte) → Aïn Draham",
        "season": "Cork stripping & malouf festival (summer)",
        "story_count": "4 stories, 1 northern loop",
    },
]

# Signature story (runs across the project) + faith-and-memory strand, for the note
SIGNATURE = (
    "Story 01 — The Charfia of Kerkennah (autumn net-setting) runs across the project as its "
    "recurring image; Story 03 (El Ghriba, spring pilgrimage) and Story 07 (the Stambeli nights) "
    "anchor a faith-and-memory strand."
)


def build_pptx(output_path: str = "INLAND_Tunisia_Pitch.pptx"):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    import os

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG = RGBColor(0xF9, 0xF4, 0xEE)
    FG = RGBColor(0x00, 0x00, 0x00)
    ACCENT = RGBColor(0x40, 0x05, 0x05)
    MUTED = RGBColor(0x5C, 0x5C, 0x5C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x0E, 0x0E, 0x0E)
    MUTED_WHITE = RGBColor(0x99, 0x99, 0x99)
    REF_COLOR = RGBColor(0x85, 0x7F, 0x76)

    HEADING = "Baskerville"
    BODY = "Baskerville"
    LABEL = "Fraunces"

    script_dir = os.path.dirname(os.path.abspath(__file__))
    logo_black = os.path.join(script_dir, "inland-logo-black.png")
    logo_white = os.path.join(script_dir, "inland-logo-white.png")

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def text(slide, left, top, w, h, txt, font, size, color,
             bold=False, align=PP_ALIGN.LEFT, italic=False, spacing=None):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = txt
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = align
        p.line_spacing = Pt(spacing) if spacing else Pt(int(size * 1.6))
        return box

    def line(slide, left, top, w, color=FG):
        shape = slide.shapes.add_connector(1, Inches(left), Inches(top), Inches(left + w), Inches(top))
        shape.line.color.rgb = color
        shape.line.width = Pt(0.75)

    def vline(slide, left, top, h, color=ACCENT):
        shape = slide.shapes.add_connector(1, Inches(left), Inches(top), Inches(left), Inches(top + h))
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

    def add_tag_pill(slide, left, top, tag_text, font=LABEL, size=9, color=MUTED):
        w = len(tag_text) * 0.075 + 0.28
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(0.3))
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        spPr = box._element.find(qn('a:spPr'))
        if spPr is None:
            spPr = box._element.makeelement(qn('a:spPr'), {})
            box._element.append(spPr)
        ln = spPr.makeelement(qn('a:ln'), {'w': '6350'})
        solidFill = ln.makeelement(qn('a:solidFill'), {})
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': 'B0A89E'})
        solidFill.append(srgb)
        ln.append(solidFill)
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'solid'}))
        spPr.append(ln)
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
        else:
            prstGeom = spPr.makeelement(qn('a:prstGeom'), {'prst': 'roundRect'})
            prstGeom.append(prstGeom.makeelement(qn('a:avLst'), {}))
            spPr.insert(0, prstGeom)
        return w

    # Title
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    logo_w, logo_h = Inches(3.5), Inches(0.735)
    logo_left = (prs.slide_width - logo_w) // 2
    if os.path.exists(logo_black):
        sl.shapes.add_picture(logo_black, logo_left, Inches(1.8), logo_w, logo_h)
    else:
        text(sl, 1, 1.6, 11.3, 0.8, "INLAND", LABEL, 48, FG, bold=True, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.9, 11.3, 1.4, "Tunisia", HEADING, 72, FG, align=PP_ALIGN.CENTER)
    line(sl, 5.5, 4.6, 2.3)
    text(sl, 1, 5.0, 11.3, 0.6, "A Mediterranean Edition", BODY, 20, MUTED, align=PP_ALIGN.CENTER)
    text(sl, 1, 6.2, 11.3, 0.4, "12 Selected Stories  ·  For the Institut français Mediterranean Photography Support  ·  2026",
         LABEL, 11, MUTED, align=PP_ALIGN.CENTER)

    # Contents
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "CONTENTS", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Overview", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    for j, (num, title) in enumerate([
        ("01", "About INLAND"), ("02", "Tunisia — What Survived"),
        ("03", "Our Angle & Categories"), ("04", "The 12 Storylines"),
        ("05", "Production Clusters"), ("06", "The Edition")]):
        y = 2.85 + j * 0.62
        text(sl, 0.8, y, 1.0, 0.5, num, HEADING, 24, ACCENT)
        text(sl, 2.0, y + 0.03, 8, 0.5, title, BODY, 20, FG)
        if j < 5:
            line(sl, 0.8, y + 0.52, 11.5, RGBColor(0xD8, 0xD2, 0xCA))

    # About
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "ABOUT", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "About INLAND", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.0, INTRO["about"], BODY, 17, FG, spacing=27)
    text(sl, 0.8, 5.3, 7.5, 1.5, INTRO["about_method"], BODY, 14, MUTED)
    text(sl, 9.5, 2.9, 3.2, 1.5, "Intimacy with\ncustodians, not\nthematic surveys.", HEADING, 24, ACCENT, italic=True)

    # Research
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THE RESEARCH", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Tunisia — What Survived", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.4, INTRO["research"], BODY, 16, FG, spacing=25)
    text(sl, 0.8, 5.9, 7.5, 1.4, INTRO["research_scope"], BODY, 13, MUTED)
    text(sl, 9.7, 2.7, 3, 1.2, "17", HEADING, 80, ACCENT)
    text(sl, 9.7, 3.8, 3, 0.4, "story leads researched", LABEL, 11, MUTED)
    text(sl, 9.7, 4.5, 3, 1.2, "12", HEADING, 80, ACCENT)
    text(sl, 9.7, 5.6, 3, 0.4, "selected for this edition", LABEL, 11, MUTED)

    # Angle & Categories
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "EDITORIAL ANGLE", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Our Angle", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.4, INTRO["angle"], BODY, 16, FG, spacing=25)
    text(sl, 0.8, 5.8, 7.5, 1.4, INTRO["angle_why"], BODY, 13, MUTED)
    text(sl, 9.7, 2.7, 3.5, 0.4, "CATEGORIES", LABEL, 11, MUTED)
    line(sl, 9.7, 3.15, 2.9)
    for j, (cat, count) in enumerate(INTRO["categories"]):
        y = 3.32 + j * 0.46
        text(sl, 9.7, y, 2.5, 0.4, cat, BODY, 14, FG)
        text(sl, 12.1, y, 1, 0.4, str(count), BODY, 14, ACCENT, align=PP_ALIGN.RIGHT, bold=True)
        if j < len(INTRO["categories"]) - 1:
            line(sl, 9.7, y + 0.39, 2.9, RGBColor(0xD8, 0xD2, 0xCA))

    # Breaker
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.0, 11.3, 0.4, "PART TWO", LABEL, 11, MUTED_WHITE, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.8, 11.3, 1.2, "The Storylines", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    text(sl, 2.5, 4.6, 8.3, 1.2,
         "Twelve selected stories across the shore and the islands — craft, "
         "ritual, fishing, food, language and landscape. Each centres a custodian, "
         "a specific place, and a practice that can still be witnessed.",
         BODY, 18, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # Story slides
    PILL = RGBColor(0x70, 0x70, 0x70)
    for i, s in enumerate(STORIES):
        sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
        text(sl, 11.0, 0.4, 2, 1.0, f"{i+1:02d}", HEADING, 64, ACCENT, align=PP_ALIGN.RIGHT)
        text(sl, 0.8, 0.6, 5, 0.4, s.category.upper(), LABEL, 11, MUTED)
        text(sl, 0.8, 1.6, 10, 1.2, s.headline, HEADING, 38, FG)
        text(sl, 0.8, 2.8, 10, 0.8, s.subtitle, HEADING, 24, MUTED, italic=True)
        line(sl, 0.8, 3.7, 11.5)
        text(sl, 0.8, 3.9, 2, 0.25, "THE STORY", LABEL, 8, PILL)
        text(sl, 0.8, 4.15, 5.8, 1.8, s.hook, BODY, 14, FG, spacing=22)
        vline(sl, 7.0, 3.9, 2.4, RGBColor(0xD8, 0xD2, 0xCA))
        text(sl, 7.4, 3.9, 2, 0.25, "OUR PROPOSAL", LABEL, 8, PILL)
        text(sl, 7.4, 4.15, 5.2, 1.8, s.proposal, BODY, 14, FG, spacing=22)
        text(sl, 0.8, 6.55, 7.2, 0.35, f"Places: {s.places}", LABEL, 10, MUTED)
        text(sl, 0.8, 6.92, 7.2, 0.3, f"Sources: {s.reference}", LABEL, 9, REF_COLOR)
        tag_x = 12.5
        for tag in reversed(s.tags):
            w = add_tag_pill(sl, 0, 6.85, tag)
            tag_x -= (w + 0.15)
            sl.shapes[-1].left = Inches(tag_x + 0.15)

    # Tracks breaker
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.0, 11.3, 0.4, "OUR SUGGESTION", LABEL, 11, MUTED_WHITE, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.8, 11.3, 1.2, "Production Clusters", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    text(sl, 2.5, 4.6, 8.3, 1.2,
         "From the twelve storylines we propose three field clusters — each "
         "grouping stories that share geography and season into one deployable "
         "shoot along Tunisia’s Mediterranean shores.", BODY, 18, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # Tracks overview
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THREE CLUSTERS", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Production Clusters", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    cw, gap, sx, top = 3.7, 0.35, 0.8, 2.9
    for j, track in enumerate(TRACKS):
        x = sx + j * (cw + gap)
        tc = RGBColor(*track["color"])
        bar = sl.shapes.add_connector(1, Inches(x), Inches(top), Inches(x + cw), Inches(top))
        bar.line.color.rgb = tc; bar.line.width = Pt(3)
        text(sl, x, top + 0.15, cw, 0.5, track["name"], HEADING, 22, tc)
        text(sl, x, top + 0.62, cw, 0.3, track["stories"], LABEL, 9, MUTED)
        text(sl, x, top + 0.95, cw, 1.7, track["desc"], BODY, 13, FG, spacing=18)
        meta = f"Geography: {track['geography']}\nSeason: {track['season']}\n{track['story_count']}"
        text(sl, x, top + 2.9, cw, 1.0, meta, LABEL, 8, REF_COLOR)
    text(sl, 0.8, 6.7, 11.7, 0.6,
         "Signature: The Charfia of Kerkennah (autumn net-setting) runs across the project; "
         "El Ghriba (spring) and the Stambeli nights anchor a faith-and-memory strand.",
         LABEL, 9, REF_COLOR)

    # Edition slide (the two IF strands)
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THE EDITION", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 11, 1.0, "Two Strands, One Edition", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    pairs = [
        ("The publication", "INLAND — yearly publication (Italy / EU). A finished Tunisia edition (~170–240 pp) on the Mediterranean theme, in the register of Issue 03, Isole Minori. The Institut français publication aid covers editing and printing in full."),
        ("The production", "The field shoot across three coastal-and-island clusters — Kerkennah, Cap Bon–Sahel, the green north — funded by the production strand for new Mediterranean photographic work."),
        ("Why the Mediterranean", "Tunisia is a Mediterranean palimpsest — Phoenician, Roman, Amazigh, Andalusian, Jewish, Ibadi and sub-Saharan — read through living custodians on its shores and islands, exactly the ‘monde méditerranéen’ the programme funds."),
        ("The principle", "Named custodians, locatable places, seasonal scenes — made with the people who keep these traditions, not about them. The loud Tunisia of revolution and migration stays as context; the quieter country is the subject."),
    ]
    for j, (h, b) in enumerate(pairs):
        y = 2.9 + j * 1.05
        text(sl, 0.8, y, 3.0, 0.4, h, LABEL, 12, ACCENT, bold=True)
        text(sl, 4.0, y - 0.02, 8.6, 1.0, b, BODY, 13, FG, spacing=18)

    # Closing
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.2, 11.3, 1.5, "12 Storylines.\nOne Tunisia.", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    lw, lh = Inches(2), Inches(0.42)
    if os.path.exists(logo_white):
        sl.shapes.add_picture(logo_white, (prs.slide_width - lw) // 2, Inches(5.0), lw, lh)
        text(sl, 1, 5.6, 11.3, 0.4, "A Mediterranean photographic edition by INLAND", LABEL, 14, MUTED_WHITE, align=PP_ALIGN.CENTER)
    else:
        text(sl, 1, 4.9, 11.3, 0.6, "A Mediterranean photographic edition by INLAND", LABEL, 14, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # strip shadows
    for sl in prs.slides:
        for shape in sl.shapes:
            spPr = shape._element.find(qn('p:spPr'))
            if spPr is not None:
                for ef in spPr.findall(qn('a:effectLst')):
                    spPr.remove(ef)
                spPr.append(spPr.makeelement(qn('a:effectLst'), {}))

    full = os.path.join(script_dir, output_path)
    prs.save(full)
    print(f"Saved PPTX: {full} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build_pptx()

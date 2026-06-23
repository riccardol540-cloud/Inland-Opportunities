"""
INLAND Angola — Presentation Builder
Generates the .pptx pitch for the Goethe-Institut Africa–Europe Partnerships
for Culture application (Visual Arts / Mobility). Mirrors the INLAND house
style. Run from this directory: `python3 build_presentation.py`
"""

from dataclasses import dataclass


@dataclass
class Story:
    headline: str
    subtitle: str
    hook: str
    proposal: str
    category: str
    tags: list
    pull_quote: str
    places: str
    reference: str


INTRO = {
    "about": (
        "INLAND is a yearly publication documenting the soul of countries "
        "before they change. Each issue embeds with named, locatable custodians "
        "of living traditions — artisans, pastoralists, cooks, singers, border "
        "families — whose stories have not yet been narrated internationally. "
        "Previous issues have covered the Faroe Islands, Georgia and Italy’s "
        "minor islands."
    ),
    "about_method": (
        "The editorial method is slow, place-based and built on trust. We spend "
        "months in pre-production research, identify specific custodians and "
        "communities, then commission photographers, writers and filmmakers to "
        "embed with them over extended field periods. The result is long-form "
        "visual storytelling grounded in named people and real places."
    ),
    "research": (
        "Angola is internationally legible through two stories — the 1975–2002 "
        "civil war and the oil-and-diamond economy that made Luanda one of the "
        "world’s most expensive cities. Both are loud; both crowd out the "
        "quieter country. INLAND asks one question across every lead: what "
        "survived? What persists in specific families, villages and practices "
        "that has been absorbed by neither the colonial inheritance, the "
        "Catholic mission, the long war, nor the petro-modernisation of the "
        "coast — and can still be witnessed and photographed today."
    ),
    "research_scope": (
        "Our desk research identified 17 story leads spanning art, ritual, "
        "governance, food, pastoral life and memory. From these, twelve "
        "storylines have been selected for editorial development — each centred "
        "on a practice that can still be witnessed, a specific place, and a "
        "named custodian or custodian community."
    ),
    "angle": (
        "We focus on endangered knowledge systems and under-documented living "
        "practices surviving despite rupture, standardisation and generational "
        "drift — from the UNESCO-listed but fragile sona sand-writing of the "
        "Lunda-Chokwe to the Kianda sea-offerings of Luanda’s original fishing "
        "people. The civil war is context, never the lens; the state is never "
        "cast as villain. Every storyline centres a custodian, a place, and a "
        "practice that can still be witnessed."
    ),
    "angle_why": (
        "This is the equitable, co-created Africa–Europe cultural work the "
        "Goethe programme exists to fund: an Angolan cultural partner and a "
        "European publication building a long-form visual record of living "
        "Angolan traditions — made with custodians, not about them, and "
        "presented in both Angola and Europe."
    ),
    "categories": [
        ("Arts & Culture", 3),
        ("History", 3),
        ("Food", 2),
        ("Myth & Religion", 2),
        ("Nature", 1),
        ("Geopolitics", 1),
    ],
}

STORIES = [
    Story(
        headline="The Men Who Write in Sand",
        subtitle="Sona, the Lunda-Chokwe Script",
        hook=(
            "In the cleaned, dotted sand of the Lunda, elder masters called akwa "
            "kuta sona draw a single unbroken geometric line that holds a proverb "
            "or a cosmology — a tradition perhaps over two thousand years old. "
            "UNESCO inscribed it as Intangible Cultural Heritage in December 2023."
        ),
        proposal=(
            "INLAND follows one master through a night of storytelling around "
            "Saurimo and Dundo, with access via Lueji A’Nkonde University — "
            "documenting a writing-in-sand recognised at the very moment it is "
            "most fragile. The defining image of the edition."
        ),
        category="Arts & Culture",
        tags=["Endangered Knowledge", "UNESCO"],
        pull_quote="A whole cosmology in one unbroken line.",
        places="Saurimo · Dundo · Lunda",
        reference="UNESCO ICH 2023; Lusona (Wikipedia); FLM Journal",
    ),
    Story(
        headline="The Fire That Stays Lit",
        subtitle="The Ovimbundu Ondjango",
        hook=(
            "The ondjango council house is built around a fire kept permanently "
            "alight — a sign of the soba’s lineage and continuity. Bailundo "
            "(Huambo) is the centre of Angola’s largest people; King Ekuikui VI, "
            "enthroned by the ombala elders in 2021, is a living authority."
        ),
        proposal=(
            "INLAND sits in on a council night and traces the wood the young "
            "bring to feed the fire — a living institution of speech and "
            "consensus, not a museum piece. The hearth is the edition’s emblem."
        ),
        category="History",
        tags=["Governance", "Continuity"],
        pull_quote="A hearth that has never been allowed to go out.",
        places="Bailundo · Huambo",
        reference="Bailundo (Wikipedia); UNILAB; Africanews",
    ),
    Story(
        headline="Hair, Herd and Drought",
        subtitle="The Mwila and the Mucubal",
        hook=(
            "On the Huíla plateau and in Namibe, the Mwila wear ochre nontombi "
            "plaits and the Mucubal the ompota headdress — hairstyles that read "
            "as a language of life-stage and mourning. The cattle economy is the "
            "base of identity, and its herd-healing knowledge is endangered."
        ),
        proposal=(
            "INLAND walks a leg of the dry-season transhumance and documents a "
            "hairstyle being made and read — a culture worn on the body and "
            "carried in the herd, under climate pressure. Custodian-led, never "
            "the tourist ‘tribal portrait’."
        ),
        category="Nature",
        tags=["Pastoralism", "Climate"],
        pull_quote="The archive is worn in the hair and walked with the cattle.",
        places="Huíla plateau · Namibe",
        reference="Minority Rights Group; PMC ethnoveterinary study 2024",
    ),
    Story(
        headline="Offerings to Kianda",
        subtitle="The Festa da Ilha",
        hook=(
            "Every second Friday of November the Axiluanda — Luanda’s original "
            "fishing people, whose name is the city’s own — cast offerings into "
            "the sea to Kianda, the water spirit who protects fishermen, as a "
            "luxury waterfront pushes them off their own sandbar."
        ),
        proposal=(
            "INLAND shoots the full midnight feast and portraits the fishers, "
            "pairing the rite with the redevelopment squeezing it out — two "
            "relationships to the same sea. The project’s signature recurring "
            "scene."
        ),
        category="Myth & Religion",
        tags=["Ritual", "Coast"],
        pull_quote="The city forgot her name; the fishermen never did.",
        places="Ilha de Luanda",
        reference="Kianda (Wikipedia); Ilha de Luanda (Wikipedia)",
    ),
    Story(
        headline="The Mask That Dances a Woman",
        subtitle="Mukanda and Mwana Pwo",
        hook=(
            "The world’s museums are full of Chokwe masks; the rite that animates "
            "them is barely filmed. During mukanda, the months-long boys’ "
            "initiation camp, a male dancer becomes mwana pwo — an idealised "
            "female ancestor — for the mothers separated from their sons."
        ),
        proposal=(
            "INLAND documents a mwana pwo performance and a carver, building a "
            "diptych of the museum vitrine against the dust of the camp. "
            "Discretion and consent around the seclusion itself."
        ),
        category="Arts & Culture",
        tags=["Initiation", "Masks"],
        pull_quote="The mask is in the museum; the dance is in the dust.",
        places="Lunda · Moxico · Dundo Museum",
        reference="Chokwe people (Wikipedia); High Museum / DIA",
    ),
    Story(
        headline="The Coffee That Went Wild",
        subtitle="Amboim Robusta",
        hook=(
            "Angola was the world’s third-largest coffee producer until the war "
            "turned plantations into wild bush. Now families harvest from trees "
            "their grandparents planted, and the EU/AFD-financed Mukafe scheme is "
            "replanting Amboim’s prized robusta in Kwanza Sul."
        ),
        proposal=(
            "INLAND follows one family through a harvest — the abandoned colonial "
            "roça now forest, the replant beside it — a story of return on land "
            "the war made wild. The family is named in the field."
        ),
        category="Food",
        tags=["Return", "Land"],
        pull_quote="They pick from trees the war turned back into forest.",
        places="Amboim · Gabela · Kwanza Sul · Uíge",
        reference="Perfect Daily Grind; STiR; Coffee in Angola (Wikipedia)",
    ),
    Story(
        headline="The Council That Outlived the King",
        subtitle="The Lumbu of Mbanza Kongo",
        hook=(
            "At Mbanza Kongo, UNESCO-listed capital of the old Kongo kingdom, "
            "there has been no king for decades — yet the lumbu, a tribunal of "
            "twenty-one elders that advised the throne for over five hundred "
            "years, still meets."
        ),
        proposal=(
            "INLAND attends a council sitting and maps the sacred topography of a "
            "city where the royal past lies underfoot — an institution that "
            "outlived the institution it served."
        ),
        category="History",
        tags=["Memory", "Kingdoms"],
        pull_quote="The throne is empty; the council still sits.",
        places="Mbanza Kongo · Zaire",
        reference="UNESCO WHC 2017; Met Museum (Kongo)",
    ),
    Story(
        headline="Njinga’s Footprints",
        subtitle="The Black Rocks of Pungo Andongo",
        hook=(
            "At the black monoliths of Pungo Andongo — the Ndongo capital and "
            "Queen Njinga’s 17th-century stronghold against the Portuguese — "
            "footprints worn into the rock are venerated locally as the queen’s "
            "own."
        ),
        proposal=(
            "INLAND shoots the footprints as a living place of veneration and the "
            "rocks as refuge — national memory made physical, a 400-year-old "
            "resistance still touched by hand."
        ),
        category="History",
        tags=["Memory", "Resistance"],
        pull_quote="A shrine without a temple.",
        places="Pungo Andongo · Malanje",
        reference="Black Rocks at Pungo Andongo (Wikipedia); Atlas Obscura",
    ),
    Story(
        headline="Those Who Ate the Ox",
        subtitle="Efundula in Cunene",
        hook=(
            "The efundula female initiation marks a girl’s passage to womanhood "
            "among the Kwanyama and Handa of Cunene — historically the symbolic "
            "‘eating of the ox’ — surviving in tension with the school and the "
            "state."
        ),
        proposal=(
            "INLAND follows the cycle with the community’s lead and the voices of "
            "young mothers caught between the rite and schooling — continuity and "
            "its costs, consent-led, never staged."
        ),
        category="Myth & Religion",
        tags=["Initiation", "Women"],
        pull_quote="To become a woman, you first ate the ox.",
        places="Cunene · southern Huíla",
        reference="Academia.edu; UCP study (Cunene)",
    ),
    Story(
        headline="Salt and the Long Boats",
        subtitle="Baía Farta",
        hook=(
            "In one Benguela bay, two squeezed coastal trades: solar salt "
            "evaporated in the pans that named the beach, and an artisanal "
            "fishing economy of tens of thousands — both pressed by industrial "
            "trawling and packaged substitutes."
        ),
        proposal=(
            "INLAND shoots a salt-raking day and a fishing dawn at the Casa dos "
            "Pescadores Artesanais — the food infrastructure of the coast, still "
            "made by hand."
        ),
        category="Food",
        tags=["Coast", "Craft"],
        pull_quote="Salt by sun, fish by hand, on one shore.",
        places="Baía Farta · Benguela",
        reference="Journey Gourmet; trackstick; DLIST Benguela",
    ),
    Story(
        headline="Clay and Fibre",
        subtitle="Women Potters and Weavers",
        hook=(
            "Coil pottery and sisal and palm basketry, made by methods passed "
            "down generations across the highlands and east, now half-"
            "commercialised for an urban market — the domestic vessels of "
            "Angolan household life, and the women who still make them."
        ),
        proposal=(
            "INLAND portraits one named maker through a full firing or weaving "
            "and the market stall beyond — women’s craft as both heritage and "
            "household economy."
        ),
        category="Arts & Culture",
        tags=["Craft", "Women"],
        pull_quote="The pot, the basket, and the hands that still make them.",
        places="Huambo · Bié · Moxico",
        reference="African Sahara; Design Milk",
    ),
    Story(
        headline="The Most Expensive City and the Musseque",
        subtitle="Two Relationships to One Shore",
        hook=(
            "Luanda’s oil-built towers stand a few hundred metres from self-built "
            "musseques and the kitanda markets that feed the city — and from the "
            "sandbar where families still fish and pray to Kianda."
        ),
        proposal=(
            "INLAND anchors the petro-city in one human — an Ilha fisher, a "
            "market woman — so the contradiction is portrait, not abstraction: "
            "two economies sharing one shore."
        ),
        category="Geopolitics",
        tags=["Modernisation", "Inequality"],
        pull_quote="Two economies, one shore.",
        places="Luanda · Ilha de Luanda",
        reference="Ilha de Luanda (Wikipedia); INLAND field framing",
    ),
]

TRACKS = [
    {
        "name": "The Lunda Cluster",
        "color": (0x2C, 0x5F, 0x2D),
        "stories": "Stories 01 + 05",
        "subtitle": "Sona & the mask rite",
        "desc": (
            "The eastern art heartland: the sona sand-writing masters around "
            "Saurimo and Dundo, and the mukanda / mwana pwo mask rite. One "
            "deployment, access via Lueji A’Nkonde University and the Dundo "
            "Museum."
        ),
        "geography": "Saurimo → Dundo (Lunda) → Moxico",
        "season": "Dry-season camps (≈May–Sep)",
        "story_count": "2 stories, 1 field base",
    },
    {
        "name": "The Planalto Cluster",
        "color": (0xB8, 0x50, 0x42),
        "stories": "Stories 02 + 06 + 11",
        "subtitle": "Council, coffee & craft",
        "desc": (
            "The central highlands: the ondjango council fire at Bailundo (King "
            "Ekuikui VI’s court), the Amboim robusta revival in Kwanza Sul, and "
            "highland potters and weavers. Governance, return and craft."
        ),
        "geography": "Bailundo (Huambo) → Amboim (Kwanza Sul)",
        "season": "Coffee harvest (≈May–Aug)",
        "story_count": "3 stories, 1 deployment",
    },
    {
        "name": "South & Coast",
        "color": (0x5B, 0x2C, 0x6F),
        "stories": "Stories 03 + 09 + 10",
        "subtitle": "Body, herd & sea",
        "desc": (
            "Mwila and Mucubal pastoral knowledge on the Huíla plateau and "
            "Namibe, the efundula female initiation in Cunene (consent-led), and "
            "salt and artisanal fishing at Baía Farta. Body, herd and sea."
        ),
        "geography": "Lubango → Cunene → Benguela coast",
        "season": "Dry-season transhumance",
        "story_count": "3 stories, seasonal windows",
    },
]

# Signature story (runs across the project) + memory strand, for the partnership note
SIGNATURE = "Story 04 — Offerings to Kianda (Festa da Ilha, Nov) runs across the project; the northern memory leads (07 Mbanza Kongo lumbu, 08 Njinga’s footprints) form a Kingdoms-and-Memory strand."


def build_pptx(output_path: str = "INLAND_Angola_Pitch.pptx"):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    import os

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG = RGBColor(0xF9, 0xF4, 0xEE)
    FG = RGBColor(0x00, 0x00, 0x00)
    ACCENT = RGBColor(0x40, 0x05, 0x05)
    MUTED = RGBColor(0x5C, 0x5C, 0x5C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x0E, 0x0E, 0x0E)
    MUTED_WHITE = RGBColor(0x99, 0x99, 0x99)
    REF_COLOR = RGBColor(0x85, 0x7F, 0x76)

    HEADING = "Baskerville"
    BODY = "Baskerville"
    LABEL = "Fraunces"

    script_dir = os.path.dirname(os.path.abspath(__file__))
    logo_black = os.path.join(script_dir, "inland-logo-black.png")
    logo_white = os.path.join(script_dir, "inland-logo-white.png")

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def text(slide, left, top, w, h, txt, font, size, color,
             bold=False, align=PP_ALIGN.LEFT, italic=False, spacing=None):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = txt
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = align
        p.line_spacing = Pt(spacing) if spacing else Pt(int(size * 1.6))
        return box

    def line(slide, left, top, w, color=FG):
        shape = slide.shapes.add_connector(1, Inches(left), Inches(top), Inches(left + w), Inches(top))
        shape.line.color.rgb = color
        shape.line.width = Pt(0.75)

    def vline(slide, left, top, h, color=ACCENT):
        shape = slide.shapes.add_connector(1, Inches(left), Inches(top), Inches(left), Inches(top + h))
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

    def add_tag_pill(slide, left, top, tag_text, font=LABEL, size=9, color=MUTED):
        w = len(tag_text) * 0.075 + 0.28
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(0.3))
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        spPr = box._element.find(qn('a:spPr'))
        if spPr is None:
            spPr = box._element.makeelement(qn('a:spPr'), {})
            box._element.append(spPr)
        ln = spPr.makeelement(qn('a:ln'), {'w': '6350'})
        solidFill = ln.makeelement(qn('a:solidFill'), {})
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': 'B0A89E'})
        solidFill.append(srgb)
        ln.append(solidFill)
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'solid'}))
        spPr.append(ln)
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
        else:
            prstGeom = spPr.makeelement(qn('a:prstGeom'), {'prst': 'roundRect'})
            prstGeom.append(prstGeom.makeelement(qn('a:avLst'), {}))
            spPr.insert(0, prstGeom)
        return w

    # Title
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    logo_w, logo_h = Inches(3.5), Inches(0.735)
    logo_left = (prs.slide_width - logo_w) // 2
    if os.path.exists(logo_black):
        sl.shapes.add_picture(logo_black, logo_left, Inches(1.8), logo_w, logo_h)
    else:
        text(sl, 1, 1.6, 11.3, 0.8, "INLAND", LABEL, 48, FG, bold=True, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.9, 11.3, 1.4, "Angola", HEADING, 72, FG, align=PP_ALIGN.CENTER)
    line(sl, 5.5, 4.6, 2.3)
    text(sl, 1, 5.0, 11.3, 0.6, "An Africa–Europe Cultural Co-Production", BODY, 20, MUTED, align=PP_ALIGN.CENTER)
    text(sl, 1, 6.2, 11.3, 0.4, "12 Selected Stories  ·  For the Goethe-Institut Africa–Europe Partnerships for Culture  ·  2026",
         LABEL, 11, MUTED, align=PP_ALIGN.CENTER)

    # Contents
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "CONTENTS", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Overview", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    for j, (num, title) in enumerate([
        ("01", "About INLAND"), ("02", "Angola — What Survived"),
        ("03", "Our Angle & Categories"), ("04", "The 12 Storylines"),
        ("05", "Production Tracks"), ("06", "The Partnership")]):
        y = 2.85 + j * 0.62
        text(sl, 0.8, y, 1.0, 0.5, num, HEADING, 24, ACCENT)
        text(sl, 2.0, y + 0.03, 8, 0.5, title, BODY, 20, FG)
        if j < 5:
            line(sl, 0.8, y + 0.52, 11.5, RGBColor(0xD8, 0xD2, 0xCA))

    # About
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "ABOUT", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "About INLAND", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.0, INTRO["about"], BODY, 17, FG, spacing=27)
    text(sl, 0.8, 5.3, 7.5, 1.5, INTRO["about_method"], BODY, 14, MUTED)
    text(sl, 9.5, 2.9, 3.2, 1.5, "Intimacy with\ncustodians, not\nthematic surveys.", HEADING, 24, ACCENT, italic=True)

    # Research
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THE RESEARCH", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Angola — What Survived", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.0, INTRO["research"], BODY, 16, FG, spacing=25)
    text(sl, 0.8, 5.7, 7.5, 1.5, INTRO["research_scope"], BODY, 13, MUTED)
    text(sl, 9.7, 2.7, 3, 1.2, "17", HEADING, 80, ACCENT)
    text(sl, 9.7, 3.8, 3, 0.4, "story leads researched", LABEL, 11, MUTED)
    text(sl, 9.7, 4.5, 3, 1.2, "12", HEADING, 80, ACCENT)
    text(sl, 9.7, 5.6, 3, 0.4, "selected for this round", LABEL, 11, MUTED)

    # Angle & Categories
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "EDITORIAL ANGLE", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Our Angle", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.2, INTRO["angle"], BODY, 16, FG, spacing=25)
    text(sl, 0.8, 5.6, 7.5, 1.5, INTRO["angle_why"], BODY, 13, MUTED)
    text(sl, 9.7, 2.7, 3.5, 0.4, "CATEGORIES", LABEL, 11, MUTED)
    line(sl, 9.7, 3.15, 2.9)
    for j, (cat, count) in enumerate(INTRO["categories"]):
        y = 3.35 + j * 0.5
        text(sl, 9.7, y, 2.5, 0.4, cat, BODY, 15, FG)
        text(sl, 12.1, y, 1, 0.4, str(count), BODY, 15, ACCENT, align=PP_ALIGN.RIGHT, bold=True)
        if j < len(INTRO["categories"]) - 1:
            line(sl, 9.7, y + 0.42, 2.9, RGBColor(0xD8, 0xD2, 0xCA))

    # Breaker
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.0, 11.3, 0.4, "PART TWO", LABEL, 11, MUTED_WHITE, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.8, 11.3, 1.2, "The Storylines", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    text(sl, 2.5, 4.6, 8.3, 1.2,
         "Twelve selected stories across art, ritual, governance, food, pastoral "
         "life and memory. Each centres a custodian, a specific place, and a "
         "practice that can still be witnessed.", BODY, 18, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # Story slides
    PILL = RGBColor(0x70, 0x70, 0x70)
    for i, s in enumerate(STORIES):
        sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
        text(sl, 11.0, 0.4, 2, 1.0, f"{i+1:02d}", HEADING, 64, ACCENT, align=PP_ALIGN.RIGHT)
        text(sl, 0.8, 0.6, 5, 0.4, s.category.upper(), LABEL, 11, MUTED)
        text(sl, 0.8, 1.6, 10, 1.2, s.headline, HEADING, 38, FG)
        text(sl, 0.8, 2.8, 10, 0.8, s.subtitle, HEADING, 24, MUTED, italic=True)
        line(sl, 0.8, 3.7, 11.5)
        text(sl, 0.8, 3.9, 2, 0.25, "THE STORY", LABEL, 8, PILL)
        text(sl, 0.8, 4.15, 5.8, 1.8, s.hook, BODY, 14, FG, spacing=22)
        vline(sl, 7.0, 3.9, 2.4, RGBColor(0xD8, 0xD2, 0xCA))
        text(sl, 7.4, 3.9, 2, 0.25, "OUR PROPOSAL", LABEL, 8, PILL)
        text(sl, 7.4, 4.15, 5.2, 1.8, s.proposal, BODY, 14, FG, spacing=22)
        text(sl, 0.8, 6.55, 7.2, 0.35, f"Places: {s.places}", LABEL, 10, MUTED)
        text(sl, 0.8, 6.92, 7.2, 0.3, f"Sources: {s.reference}", LABEL, 9, REF_COLOR)
        tag_x = 12.5
        for tag in reversed(s.tags):
            w = add_tag_pill(sl, 0, 6.85, tag)
            tag_x -= (w + 0.15)
            sl.shapes[-1].left = Inches(tag_x + 0.15)

    # Tracks breaker
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.0, 11.3, 0.4, "OUR SUGGESTION", LABEL, 11, MUTED_WHITE, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.8, 11.3, 1.2, "Production Tracks", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    text(sl, 2.5, 4.6, 8.3, 1.2,
         "From the twelve storylines we propose three production tracks for a "
         "pilot — each grouping stories that share geography and field timing "
         "into one deployable trip with the Angolan partner.", BODY, 18, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # Tracks overview
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THREE TRACKS", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Production Tracks", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    cw, gap, sx, top = 3.7, 0.35, 0.8, 2.9
    for j, track in enumerate(TRACKS):
        x = sx + j * (cw + gap)
        tc = RGBColor(*track["color"])
        bar = sl.shapes.add_connector(1, Inches(x), Inches(top), Inches(x + cw), Inches(top))
        bar.line.color.rgb = tc; bar.line.width = Pt(3)
        text(sl, x, top + 0.15, cw, 0.5, track["name"], HEADING, 23, tc)
        text(sl, x, top + 0.62, cw, 0.3, track["stories"], LABEL, 9, MUTED)
        text(sl, x, top + 0.95, cw, 1.7, track["desc"], BODY, 13, FG, spacing=18)
        meta = f"Geography: {track['geography']}\nSeason: {track['season']}\n{track['story_count']}"
        text(sl, x, top + 2.9, cw, 1.0, meta, LABEL, 8, REF_COLOR)
    text(sl, 0.8, 6.7, 11.7, 0.6,
         "Signature: Offerings to Kianda (Festa da Ilha, Nov) runs across the project; "
         "the Mbanza Kongo lumbu and Njinga’s footprints form a Kingdoms-and-Memory strand.",
         LABEL, 9, REF_COLOR)

    # Partnership slide
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THE PARTNERSHIP", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 11, 1.0, "An Equitable Co-Production", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    pairs = [
        ("European partner", "INLAND — yearly publication (Italy / EU). Editorial, photography and publishing; presentation of the work in Europe."),
        ("African partner", "An Angolan cultural organisation (lead candidate: Lueji A’Nkonde University, Lunda Norte — which coordinated the sona UNESCO nomination) — co-creation, community access, custodian liaison, presentation in Angola."),
        ("Funding path", "Mobility grant (≤€4,000) funds the scouting trip; the Visual Arts grant funds collaborative production and a public presentation in Angola and Europe."),
        ("The principle", "Made with custodians, not about them — equitable Africa–Europe co-creation, exactly the dialogue the Goethe programme exists to fund."),
    ]
    for j, (h, b) in enumerate(pairs):
        y = 2.9 + j * 1.05
        text(sl, 0.8, y, 3.0, 0.4, h, LABEL, 12, ACCENT, bold=True)
        text(sl, 4.0, y - 0.02, 8.6, 1.0, b, BODY, 14, FG, spacing=19)

    # Closing
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.2, 11.3, 1.5, "12 Storylines.\nOne Angola.", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    lw, lh = Inches(2), Inches(0.42)
    if os.path.exists(logo_white):
        sl.shapes.add_picture(logo_white, (prs.slide_width - lw) // 2, Inches(5.0), lw, lh)
        text(sl, 1, 5.6, 11.3, 0.4, "An Africa–Europe co-production proposal", LABEL, 14, MUTED_WHITE, align=PP_ALIGN.CENTER)
    else:
        text(sl, 1, 4.9, 11.3, 0.6, "An Africa–Europe co-production proposal by INLAND", LABEL, 14, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # strip shadows
    for sl in prs.slides:
        for shape in sl.shapes:
            spPr = shape._element.find(qn('p:spPr'))
            if spPr is not None:
                for ef in spPr.findall(qn('a:effectLst')):
                    spPr.remove(ef)
                spPr.append(spPr.makeelement(qn('a:effectLst'), {}))

    full = os.path.join(script_dir, output_path)
    prs.save(full)
    print(f"Saved PPTX: {full} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build_pptx()

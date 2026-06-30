"""
INLAND Georgia — Presentation Builder
Generates the .pptx pitch for the W. Eugene Smith Grant in Humanistic Photography.
Anchored on Georgia (INLAND Issue 02, 2024) — a proposed long-form humanistic essay
"The Quiet Country" among Georgia's highland custodians, in Smith's concerned tradition.
Mirrors the INLAND house style. Run from this directory: `python3 build_presentation.py`
"""

from dataclasses import dataclass


@dataclass
class Story:
    headline: str
    subtitle: str
    hook: str
    proposal: str
    category: str
    tags: list
    pull_quote: str
    places: str
    reference: str


INTRO = {
    "about": (
        "INLAND is a yearly publication documenting the soul of countries "
        "before they change. Each issue embeds with named, locatable custodians "
        "of living traditions — artisans, pastoralists, cooks, singers, border "
        "families — whose stories have not yet been narrated internationally. "
        "Previous issues have covered the Faroe Islands, Georgia and Italy’s "
        "minor islands."
    ),
    "about_method": (
        "The editorial method is slow, place-based and built on trust. We spend "
        "months in pre-production research, identify specific custodians and "
        "communities, then commission photographers, writers and filmmakers to "
        "embed with them over extended field periods. The result is long-form "
        "visual storytelling grounded in named people and real places."
    ),
    "research": (
        "Georgia is internationally legible through three loud stories — the 2008 "
        "war and the occupied lines of Abkhazia and South Ossetia, the recurring "
        "Tbilisi protests over the country’s European future, and the "
        "wine-and-feast tourism brand. All three are loud; all three crowd out the "
        "quieter country. INLAND asks one question across every lead: what "
        "survived? What persists in the high valleys and a few plains communities "
        "that has been absorbed by neither empire, Soviet rule, the market, nor "
        "the national brand — and can still be witnessed and photographed today."
    ),
    "research_scope": (
        "Our desk research identified 16 story leads across pastoral life, song, "
        "faith, craft, food and memory. From these, twelve have been selected for "
        "the proposed project — each centred on a practice that can still be "
        "witnessed, a specific place, and a named custodian or community in the "
        "tradition of W. Eugene Smith’s concerned, humanistic documentary."
    ),
    "angle": (
        "We read Georgia through the people who keep its endangered knowledge: the "
        "mounted transhumance over the Abano Pass, the defence-tower families of "
        "Ushguli, the pre-Christian sun-hymns of Svaneti, the only women’s Sufi "
        "zikr in the South Caucasus at Duisi, the wine made under the floor of the "
        "house, and a 26-century Jewish diaspora down to its last worshippers. The "
        "occupation lines and the protests are context, never the lens; the state "
        "is never cast as villain."
    ),
    "angle_why": (
        "This is the humanistic documentary the W. Eugene Smith Grant exists to "
        "fund: a long-term, compassionate essay about the human condition, made "
        "with named custodians rather than about them. INLAND’s published Georgia "
        "edition (Issue 02, 2024) is the photographer’s prior cohesive body of "
        "work; the grant would underwrite the extended field embedding that "
        "carries it deeper — continuity over crisis, dignity over spectacle."
    ),
    "categories": [
        ("Arts & Culture", 3),
        ("Food", 2),
        ("Myth & Religion", 2),
        ("Local communities", 2),
        ("History", 1),
        ("Nature", 1),
        ("Geopolitics", 1),
    ],
}

STORIES = [
    Story(
        headline="The Road in May",
        subtitle="The Transhumance over the Abano Pass",
        hook=(
            "Around a tenth of Tusheti’s people still live to the rhythm of "
            "transhumance, driving sheep — and, ever more rarely, riding the "
            "migration on horseback — between the Kakheti lowlands and the high "
            "summer grass of Tusheti, climbing and descending the Abano Pass at "
            "2,850 metres, one of the most dangerous roads in the country."
        ),
        proposal=(
            "INLAND follows one shepherd family through the spring drive up and the "
            "autumn descent — the night camps, the dogs, the pass in cloud — a "
            "half-millennium-old movement of animals and people, on a real date. "
            "The project’s signature recurring image."
        ),
        category="Nature",
        tags=["Pastoral", "Season"],
        pull_quote="A road older than the nation, still walked once a year.",
        places="Kakheti → Abano Pass → Tusheti (Omalo)",
        reference="Atlas Obscura; National Geographic; RFE/RL (horse migration); Milk Trekker",
    ),
    Story(
        headline="The Tower and the Byre",
        subtitle="The Tower-Keeping Families of Ushguli",
        hook=(
            "Ushguli, at the confluence of the Inguri and Black rivers, is among "
            "Europe’s highest inhabited places — a UNESCO site of more than 200 "
            "medieval koshki defence towers. A handful of families still keep house, "
            "byre and a standing tower together; Roland and Gulo Chelidze tend "
            "cattle, pigs and horses below the family’s keep."
        ),
        proposal=(
            "INLAND spends a winter day with one Ushguli family between the byre and "
            "the tower — the stock, the horse transport across roadless country, the "
            "tower that was once home, fortress and treasury at once."
        ),
        category="History",
        tags=["UNESCO", "Mountain"],
        pull_quote="A house, a fortress and a treasury — the same stone.",
        places="Ushguli · Chazhashi · Upper Svaneti",
        reference="UNESCO WHC 709; National Geographic; Atlas Obscura; Museum of Wander",
    ),
    Story(
        headline="Lile, the Sun Song",
        subtitle="The Animist Hymns of Svaneti",
        hook=(
            "Some of the oldest layers of Georgian polyphony are Svan: hymns like "
            "Lile — the ancient Svan word for the Sun God — and Kviria carry largely "
            "pre-Christian, animist texts that Svans today often re-read through a "
            "Christian lens. The Mestia ensemble Riho kept them for over fifty years."
        ),
        proposal=(
            "INLAND sits a rehearsal with Riho, now directed by Vakhtang Pilpani — "
            "three male voices braided into a sun-hymn older than the church, in a "
            "Mestia front room as the snow comes down."
        ),
        category="Arts & Culture",
        tags=["Music", "Endangered"],
        pull_quote="A hymn to the sun, still sung where the church put its name on it.",
        places="Mestia · Svaneti",
        reference="Georgian Chant (Kviria; Svan chant); Chanting.ge; Discogs (Riho)",
    ),
    Story(
        headline="The Highest Voice",
        subtitle="Gurian Krimanchuli on the Black Sea",
        hook=(
            "On the Black Sea side, Gurian polyphony is the most acrobatic of all "
            "Georgian song: above the texture flies the krimanchuli, a high "
            "yodel-falsetto of tension and release. In Lanchkhuti, the Turkia family "
            "has been central to keeping and teaching the Gurian repertoire."
        ),
        proposal=(
            "INLAND works a supra table in Guria where a three-part song breaks into "
            "the krimanchuli — the voice as a household inheritance, taught father to "
            "son across the feast."
        ),
        category="Arts & Culture",
        tags=["UNESCO", "Music"],
        pull_quote="The voice goes up where no other voice can follow.",
        places="Lanchkhuti · Guria",
        reference="UNESCO ICH 00008; chortownia.org; harmonychroniclesmag.com",
    ),
    Story(
        headline="The Fire in the Enamel",
        subtitle="The Minankari Revival of Tbilisi",
        hook=(
            "Minankari, Georgian cloisonné enamel, reached its height between the "
            "10th and 15th centuries on icons and crosses, then was abandoned for "
            "centuries. Over the last two decades Tbilisi artists have revived it by "
            "an almost unchanged hand process, centred on the ‘Ornament’ gallery."
        ),
        proposal=(
            "INLAND portraits a master — Thea Gurgenidze among them — laying gold "
            "wire into cells and firing colour at the bench: a medieval church "
            "technique alive in a living hand, taught to an apprentice."
        ),
        category="Arts & Culture",
        tags=["Craft", "Revival"],
        pull_quote="A church technique, abandoned for centuries, back in living hands.",
        places="Tbilisi · the ‘Ornament’ gallery",
        reference="NTD; mygeotrip (Minankari; Ornament); georgia.to",
    ),
    Story(
        headline="The Wine Under the Floor",
        subtitle="Qvevri & the Marani at Harvest",
        hook=(
            "Georgia’s 8,000 years of winemaking are lived, not marketed, in the "
            "marani — the cellar where egg-shaped qvevri are buried to the neck and "
            "the wine ferments on its skins through the cold without machinery. The "
            "qvevri themselves are still thrown by a few families at Vardisubani."
        ),
        proposal=(
            "INLAND follows a Kakhetian family through rtveli, the harvest — the "
            "pressing, the qvevri sealed underground, the feast — and the potter at "
            "Vardisubani firing a vessel taller than a man."
        ),
        category="Food",
        tags=["UNESCO", "Harvest"],
        pull_quote="The wine is made under the floor of the house, on its own skins.",
        places="Kakheti · Vardisubani · Telavi",
        reference="UNESCO ICH 00870; wine.gov.ge; Decanter; winesgeorgia.com",
    ),
    Story(
        headline="The Cheese in the Skin",
        subtitle="Tushetian Guda on the Summer Mountain",
        hook=(
            "Tushetian guda is made by herders during the summer transhumance from "
            "whole raw milk, matured inside a guda — a whole sheepskin sack — and "
            "covered with the felted nabadi burka so it takes the salt evenly. The "
            "women run the dairy and the wool; the Alaznistavi cooperative has "
            "revived the craft."
        ),
        proposal=(
            "INLAND documents a summer dairy hut — milk into skin, the nabadi over "
            "the curing room, the pardagi felt carpets — naming a guda maker with "
            "the Alaznistavi women."
        ),
        category="Food",
        tags=["Pastoral", "Women"],
        pull_quote="A cheese born in a sheepskin, on a mountain, in summer.",
        places="Tusheti · Alaznistavi cooperative",
        reference="EU4Georgia (Alaznistavi); Milk Trekker; Wikipedia (Tushetians)",
    ),
    Story(
        headline="The Mountain’s Own Faith",
        subtitle="The Shrine Cult of Khevsureti",
        hook=(
            "In Khevsureti people worship at jvari / khati — stone shrines two metres "
            "tall, hung with ibex horns and a bell — in a syncretism where, behind "
            "the Christian names and candles, the rite is essentially pre-Christian, "
            "served by an elected khevisberi, a dasturi, and a ‘hostess of the shrine.’"
        ),
        proposal=(
            "INLAND works Atengenoba, the great festival 100 days after Easter: the "
            "khevisberi at the horned shrine, the beer brewed on the spot, the "
            "offering — at the fortress villages of Shatili and Mutso. Consent-led."
        ),
        category="Myth & Religion",
        tags=["Ritual", "Pre-Christian"],
        pull_quote="Behind the candles and the cross, a faith older than the church.",
        places="Khevsureti · Shatili · Mutso",
        reference="Wikipedia (Khevsureti); academia.edu (Khatis Diasakhlisi); khconsult",
    ),
    Story(
        headline="The Women Who Sing the Zikr",
        subtitle="The Kist Sufi Circle of Pankisi",
        hook=(
            "The Kist of Pankisi descend from Chechens who crossed the mountains in "
            "the 1800s and follow the pacifist Sufi path of Kunta-Hajji. Pankisi is "
            "the only place in the South Caucasus where women lead the zikr in the "
            "mosque — every Friday at the Old Mosque in Duisi — and the elders fear "
            "it will fall silent."
        ),
        proposal=(
            "INLAND photographs the Friday women’s circle with consent — the turning "
            "chant, the nazms in the Kist tongue, the call for ‘peace in the "
            "Caucasus’ — with the Aznash Laaman singers. Ritual, not spectacle; "
            "identities protected where asked."
        ),
        category="Myth & Religion",
        tags=["Sufi", "Endangered"],
        pull_quote="The only place women lead the zikr — and it is going quiet.",
        places="Duisi · Pankisi Gorge",
        reference="OC Media; Al Jazeera America; pankisivalley.com",
    ),
    Story(
        headline="The Saint in the Desert",
        subtitle="David Gareja on the Disputed Line",
        hook=(
            "David Gareja is a complex of 19 medieval cave monasteries — some 5,000 "
            "cells cut into the Kakheti semi-desert in the 6th century, with "
            "8th–13th-century frescoes including a rare portrait of Queen Tamar. The "
            "Soviet-drawn Georgia–Azerbaijan border runs straight through it."
        ),
        proposal=(
            "INLAND frames monastic life in the rock against the quiet frontier "
            "dispute — a monk on the desert ridge where the boundary meets the "
            "painted caves. Tension, not attack; verify border access first."
        ),
        category="Geopolitics",
        tags=["Borderland", "Monastic"],
        pull_quote="A 6th-century desert, cut in two by a 20th-century line.",
        places="David Gareja · Udabno · Kakheti",
        reference="UNESCO tentative 5224; Advantour; Atlas Obscura; wander-lush.org",
    ),
    Story(
        headline="The Granary Language",
        subtitle="Svan, Dying at Home, Sung by the Young",
        hook=(
            "Svan — an unwritten Kartvelian language perhaps three millennia diverged "
            "from Georgian — is endangered, spoken by a shrinking number of elders, "
            "even as Svan music enjoys a youth-led revival driven from Mestia. The "
            "songs return to young mouths while the everyday language thins."
        ),
        proposal=(
            "INLAND follows young singers of the Lileo ensemble in Mestia — a "
            "rehearsal where teenagers sing hymns in a tongue fewer of them speak at "
            "home — the old words kept like seed-grain."
        ),
        category="Local communities",
        tags=["Language", "Youth"],
        pull_quote="They sing the words their grandparents spoke, and they no longer do.",
        places="Mestia · Ushguli · Svaneti",
        reference="Chai Khana; georgianchant.org; UNESCO Atlas of endangered languages",
    ),
    Story(
        headline="2,600 Years of Welcome",
        subtitle="The Last Worshippers of Oni & Kutaisi",
        hook=(
            "Jews have lived in Georgia for some 26 centuries, in what is often "
            "called the one country where they were never persecuted as Jews. The "
            "Oni synagogue (1895) once held over a thousand; today fewer than 20 "
            "worship there. Kutaisi’s Jewish quarter, once 30,000 strong, counts "
            "fewer than a hundred families."
        ),
        proposal=(
            "INLAND documents a Sabbath where the minyan is barely made, in a "
            "synagogue built for a thousand — the keepers of one of the world’s "
            "oldest diasporas, with discretion."
        ),
        category="Local communities",
        tags=["Diaspora", "Endangered"],
        pull_quote="A synagogue built for a thousand, kept by twenty.",
        places="Oni (Racha) · Kutaisi",
        reference="Wikipedia (Oni; Kutaisi synagogues); Gil Travel; jguideeurope.org",
    ),
]

TRACKS = [
    {
        "name": "The Eastern Highlands",
        "color": (0x1B, 0x49, 0x65),
        "stories": "Stories 01 + 07 + 08 + 09",
        "subtitle": "Transhumance, cheese, shrine & zikr",
        "desc": (
            "From a Kakheti base: the Abano-Pass transhumance and Tushetian guda "
            "cheese, the shrine cult of Khevsureti at Atengenoba, and the women’s "
            "Sufi zikr of Pankisi. One eastern deployment across the high summer."
        ),
        "geography": "Kakheti (Telavi/Akhmeta) → Tusheti → Khevsureti → Pankisi (Duisi)",
        "season": "Transhumance May & Sept; Atengenoba ~mid-Jul; guda Jun–Sep; zikr Fri",
        "story_count": "4 stories, 1 eastern base",
    },
    {
        "name": "The Western Voices",
        "color": (0x3A, 0x5A, 0x40),
        "stories": "Stories 02 + 03 + 04 + 11 + 12",
        "subtitle": "Towers, song, language & diaspora",
        "desc": (
            "The west and the high north-west: the tower families and animist "
            "sun-hymns of Svaneti, the endangered Svan language sung by the young, "
            "the Gurian krimanchuli of Lanchkhuti, and the last Jewish worshippers "
            "of Oni and Kutaisi."
        ),
        "geography": "Mestia/Ushguli (Svaneti) → Lanchkhuti (Guria) → Oni (Racha) → Kutaisi",
        "season": "Svaneti late spring–autumn; rehearsals & Sabbaths year-round",
        "story_count": "5 stories, 1 western loop",
    },
    {
        "name": "The Plain & the Desert",
        "color": (0x5B, 0x2C, 0x6F),
        "stories": "Stories 05 + 06 + 10",
        "subtitle": "Enamel, wine & the border monastery",
        "desc": (
            "Tbilisi and the Kakhetian plain: the minankari enamel revival, the "
            "qvevri marani at harvest, and the cave monastery of David Gareja on the "
            "disputed desert line. Anchored on rtveli, the autumn harvest."
        ),
        "geography": "Tbilisi → Kakheti (Vardisubani/Telavi) → David Gareja (Udabno)",
        "season": "Rtveli late Sep; enamel year-round; Gareja spring/autumn light",
        "story_count": "3 stories, autumn window",
    },
]

# Signature image (runs across the project) + binding social thread, for the note
SIGNATURE = (
    "Story 01 — The Road in May (the Abano Pass crossing) runs across the project as its "
    "recurring image; the supra and Georgian polyphony (Stories 03, 04) bind the clusters "
    "as the project’s social thread."
)


def build_pptx(output_path: str = "INLAND_Georgia_Pitch.pptx"):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    import os

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG = RGBColor(0xF9, 0xF4, 0xEE)
    FG = RGBColor(0x00, 0x00, 0x00)
    ACCENT = RGBColor(0x40, 0x05, 0x05)
    MUTED = RGBColor(0x5C, 0x5C, 0x5C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x0E, 0x0E, 0x0E)
    MUTED_WHITE = RGBColor(0x99, 0x99, 0x99)
    REF_COLOR = RGBColor(0x85, 0x7F, 0x76)

    HEADING = "Baskerville"
    BODY = "Baskerville"
    LABEL = "Fraunces"

    script_dir = os.path.dirname(os.path.abspath(__file__))
    logo_black = os.path.join(script_dir, "inland-logo-black.png")
    logo_white = os.path.join(script_dir, "inland-logo-white.png")

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def text(slide, left, top, w, h, txt, font, size, color,
             bold=False, align=PP_ALIGN.LEFT, italic=False, spacing=None):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = txt
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = align
        p.line_spacing = Pt(spacing) if spacing else Pt(int(size * 1.6))
        return box

    def line(slide, left, top, w, color=FG):
        shape = slide.shapes.add_connector(1, Inches(left), Inches(top), Inches(left + w), Inches(top))
        shape.line.color.rgb = color
        shape.line.width = Pt(0.75)

    def vline(slide, left, top, h, color=ACCENT):
        shape = slide.shapes.add_connector(1, Inches(left), Inches(top), Inches(left), Inches(top + h))
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

    def add_tag_pill(slide, left, top, tag_text, font=LABEL, size=9, color=MUTED):
        w = len(tag_text) * 0.075 + 0.28
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(0.3))
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        spPr = box._element.find(qn('a:spPr'))
        if spPr is None:
            spPr = box._element.makeelement(qn('a:spPr'), {})
            box._element.append(spPr)
        ln = spPr.makeelement(qn('a:ln'), {'w': '6350'})
        solidFill = ln.makeelement(qn('a:solidFill'), {})
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': 'B0A89E'})
        solidFill.append(srgb)
        ln.append(solidFill)
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'solid'}))
        spPr.append(ln)
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
        else:
            prstGeom = spPr.makeelement(qn('a:prstGeom'), {'prst': 'roundRect'})
            prstGeom.append(prstGeom.makeelement(qn('a:avLst'), {}))
            spPr.insert(0, prstGeom)
        return w

    # Title
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    logo_w, logo_h = Inches(3.5), Inches(0.735)
    logo_left = (prs.slide_width - logo_w) // 2
    if os.path.exists(logo_black):
        sl.shapes.add_picture(logo_black, logo_left, Inches(1.8), logo_w, logo_h)
    else:
        text(sl, 1, 1.6, 11.3, 0.8, "INLAND", LABEL, 48, FG, bold=True, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.9, 11.3, 1.4, "Georgia", HEADING, 72, FG, align=PP_ALIGN.CENTER)
    line(sl, 5.5, 4.6, 2.3)
    text(sl, 1, 5.0, 11.3, 0.6, "The Quiet Country — A Humanistic Documentary", BODY, 20, MUTED, align=PP_ALIGN.CENTER)
    text(sl, 1, 6.2, 11.3, 0.4, "12 Selected Stories  ·  For the W. Eugene Smith Grant in Humanistic Photography  ·  2026",
         LABEL, 11, MUTED, align=PP_ALIGN.CENTER)

    # Contents
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "CONTENTS", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Overview", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    for j, (num, title) in enumerate([
        ("01", "About INLAND"), ("02", "Georgia — What Survived"),
        ("03", "Our Angle & Categories"), ("04", "The 12 Storylines"),
        ("05", "Production Clusters"), ("06", "The Application")]):
        y = 2.85 + j * 0.62
        text(sl, 0.8, y, 1.0, 0.5, num, HEADING, 24, ACCENT)
        text(sl, 2.0, y + 0.03, 8, 0.5, title, BODY, 20, FG)
        if j < 5:
            line(sl, 0.8, y + 0.52, 11.5, RGBColor(0xD8, 0xD2, 0xCA))

    # About
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "ABOUT", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "About INLAND", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.0, INTRO["about"], BODY, 17, FG, spacing=27)
    text(sl, 0.8, 5.3, 7.5, 1.5, INTRO["about_method"], BODY, 14, MUTED)
    text(sl, 9.5, 2.9, 3.2, 1.5, "Intimacy with\ncustodians, not\nthematic surveys.", HEADING, 24, ACCENT, italic=True)

    # Research
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THE RESEARCH", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Georgia — What Survived", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.4, INTRO["research"], BODY, 16, FG, spacing=25)
    text(sl, 0.8, 5.9, 7.5, 1.4, INTRO["research_scope"], BODY, 13, MUTED)
    text(sl, 9.7, 2.7, 3, 1.2, "16", HEADING, 80, ACCENT)
    text(sl, 9.7, 3.8, 3, 0.4, "story leads researched", LABEL, 11, MUTED)
    text(sl, 9.7, 4.5, 3, 1.2, "12", HEADING, 80, ACCENT)
    text(sl, 9.7, 5.6, 3, 0.4, "selected for the project", LABEL, 11, MUTED)

    # Angle & Categories
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "EDITORIAL ANGLE", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Our Angle", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.4, INTRO["angle"], BODY, 16, FG, spacing=25)
    text(sl, 0.8, 5.8, 7.5, 1.4, INTRO["angle_why"], BODY, 13, MUTED)
    text(sl, 9.7, 2.7, 3.5, 0.4, "CATEGORIES", LABEL, 11, MUTED)
    line(sl, 9.7, 3.15, 2.9)
    for j, (cat, count) in enumerate(INTRO["categories"]):
        y = 3.32 + j * 0.46
        text(sl, 9.7, y, 2.5, 0.4, cat, BODY, 14, FG)
        text(sl, 12.1, y, 1, 0.4, str(count), BODY, 14, ACCENT, align=PP_ALIGN.RIGHT, bold=True)
        if j < len(INTRO["categories"]) - 1:
            line(sl, 9.7, y + 0.39, 2.9, RGBColor(0xD8, 0xD2, 0xCA))

    # Breaker
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.0, 11.3, 0.4, "PART TWO", LABEL, 11, MUTED_WHITE, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.8, 11.3, 1.2, "The Storylines", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    text(sl, 2.5, 4.6, 8.3, 1.2,
         "Twelve selected stories across pasture, song, faith, craft, food and "
         "memory. Each centres a named custodian, a specific place, and a practice "
         "that can still be witnessed — in the tradition of concerned photography.",
         BODY, 18, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # Story slides
    PILL = RGBColor(0x70, 0x70, 0x70)
    for i, s in enumerate(STORIES):
        sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
        text(sl, 11.0, 0.4, 2, 1.0, f"{i+1:02d}", HEADING, 64, ACCENT, align=PP_ALIGN.RIGHT)
        text(sl, 0.8, 0.6, 5, 0.4, s.category.upper(), LABEL, 11, MUTED)
        text(sl, 0.8, 1.6, 10, 1.2, s.headline, HEADING, 38, FG)
        text(sl, 0.8, 2.8, 10, 0.8, s.subtitle, HEADING, 24, MUTED, italic=True)
        line(sl, 0.8, 3.7, 11.5)
        text(sl, 0.8, 3.9, 2, 0.25, "THE STORY", LABEL, 8, PILL)
        text(sl, 0.8, 4.15, 5.8, 1.8, s.hook, BODY, 14, FG, spacing=22)
        vline(sl, 7.0, 3.9, 2.4, RGBColor(0xD8, 0xD2, 0xCA))
        text(sl, 7.4, 3.9, 2, 0.25, "OUR PROPOSAL", LABEL, 8, PILL)
        text(sl, 7.4, 4.15, 5.2, 1.8, s.proposal, BODY, 14, FG, spacing=22)
        text(sl, 0.8, 6.55, 7.2, 0.35, f"Places: {s.places}", LABEL, 10, MUTED)
        text(sl, 0.8, 6.92, 7.2, 0.3, f"Sources: {s.reference}", LABEL, 9, REF_COLOR)
        tag_x = 12.5
        for tag in reversed(s.tags):
            w = add_tag_pill(sl, 0, 6.85, tag)
            tag_x -= (w + 0.15)
            sl.shapes[-1].left = Inches(tag_x + 0.15)

    # Tracks breaker
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.0, 11.3, 0.4, "OUR SUGGESTION", LABEL, 11, MUTED_WHITE, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.8, 11.3, 1.2, "Production Clusters", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    text(sl, 2.5, 4.6, 8.3, 1.2,
         "From the twelve storylines we propose three field clusters — each "
         "grouping stories that share geography and season into one deployable "
         "embed across Georgia’s highlands and plain.", BODY, 18, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # Tracks overview
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THREE CLUSTERS", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Production Clusters", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    cw, gap, sx, top = 3.7, 0.35, 0.8, 2.9
    for j, track in enumerate(TRACKS):
        x = sx + j * (cw + gap)
        tc = RGBColor(*track["color"])
        bar = sl.shapes.add_connector(1, Inches(x), Inches(top), Inches(x + cw), Inches(top))
        bar.line.color.rgb = tc; bar.line.width = Pt(3)
        text(sl, x, top + 0.15, cw, 0.5, track["name"], HEADING, 22, tc)
        text(sl, x, top + 0.62, cw, 0.3, track["stories"], LABEL, 9, MUTED)
        text(sl, x, top + 0.95, cw, 1.7, track["desc"], BODY, 13, FG, spacing=18)
        meta = f"Geography: {track['geography']}\nSeason: {track['season']}\n{track['story_count']}"
        text(sl, x, top + 2.9, cw, 1.0, meta, LABEL, 8, REF_COLOR)
    text(sl, 0.8, 6.7, 11.7, 0.6,
         "Signature: The Road in May (the Abano Pass crossing) runs across the project; "
         "the supra and Georgian polyphony bind the clusters as its social thread.",
         LABEL, 9, REF_COLOR)

    # Application slide (the Smith framing)
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THE APPLICATION", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 11, 1.0, "A Humanistic Project, in Smith’s Tradition", HEADING, 36, FG)
    line(sl, 0.8, 2.5, 11.5)
    pairs = [
        ("The applicant", "An INLAND-commissioned documentary photographer is the named applicant; INLAND is the editorial home that researches, embeds and publishes the work."),
        ("Past work", "INLAND’s Georgia edition (Issue 02, 2024) is the prior cohesive body of work — long-form, humanistic, custodian-led — that this project carries deeper."),
        ("Proposed project", "‘The Quiet Country’: a year among Georgia’s highland custodians — shepherds, tower families, shrine-priests, the Pankisi zikr women, winemakers — up to 40 images."),
        ("Why Smith", "Concerned, compassionate documentary about the human condition; continuity over crisis, dignity over spectacle. The grant underwrites the extended field embedding."),
    ]
    for j, (h, b) in enumerate(pairs):
        y = 2.9 + j * 1.05
        text(sl, 0.8, y, 3.0, 0.4, h, LABEL, 12, ACCENT, bold=True)
        text(sl, 4.0, y - 0.02, 8.6, 1.0, b, BODY, 13, FG, spacing=18)

    # Closing
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.2, 11.3, 1.5, "12 Storylines.\nOne Quiet Country.", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    lw, lh = Inches(2), Inches(0.42)
    if os.path.exists(logo_white):
        sl.shapes.add_picture(logo_white, (prs.slide_width - lw) // 2, Inches(5.0), lw, lh)
        text(sl, 1, 5.6, 11.3, 0.4, "A humanistic documentary project by INLAND, in the tradition of W. Eugene Smith", LABEL, 13, MUTED_WHITE, align=PP_ALIGN.CENTER)
    else:
        text(sl, 1, 4.9, 11.3, 0.6, "A humanistic documentary project by INLAND, in the tradition of W. Eugene Smith", LABEL, 13, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # strip shadows
    for sl in prs.slides:
        for shape in sl.shapes:
            spPr = shape._element.find(qn('p:spPr'))
            if spPr is not None:
                for ef in spPr.findall(qn('a:effectLst')):
                    spPr.remove(ef)
                spPr.append(spPr.makeelement(qn('a:effectLst'), {}))

    full = os.path.join(script_dir, output_path)
    prs.save(full)
    print(f"Saved PPTX: {full} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build_pptx()

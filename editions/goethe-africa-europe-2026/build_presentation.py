"""
INLAND Rwanda — Presentation Builder
Generates the .pptx pitch for the Goethe-Institut Africa–Europe Partnerships
for Culture application (Visual Arts / Mobility). Mirrors the INLAND Kashmir
house style. Run from this directory: `python3 build_presentation.py`
"""

from dataclasses import dataclass


@dataclass
class Story:
    headline: str
    subtitle: str
    hook: str
    proposal: str
    category: str
    tags: list
    pull_quote: str
    places: str
    reference: str


INTRO = {
    "about": (
        "INLAND is a yearly publication documenting the soul of countries "
        "before they change. Each issue embeds with named, locatable custodians "
        "of living traditions — artisans, pastoralists, cooks, singers, border "
        "families — whose stories have not yet been narrated internationally. "
        "Previous issues have covered the Faroe Islands, Georgia and Italy’s "
        "minor islands."
    ),
    "about_method": (
        "The editorial method is slow, place-based and built on trust. We spend "
        "months in pre-production research, identify specific custodians and "
        "communities, then commission photographers, writers and filmmakers to "
        "embed with them over extended field periods. The result is long-form "
        "visual storytelling grounded in named people and real places."
    ),
    "research": (
        "Rwanda is internationally legible through two stories — the 1994 "
        "genocide and the development ‘miracle’. Both are loud; both crowd out "
        "the quieter country. INLAND asks one question across every lead: what "
        "survived? What persists in specific families, villages and practices "
        "that has been absorbed by neither the colonial inheritance, the "
        "Christian mission, nor the modernisation project — and can still be "
        "witnessed and photographed today."
    ),
    "research_scope": (
        "Our first round of research identified 20 story leads spanning "
        "communities, ritual, food, ecology, craft and music. From these, "
        "twelve storylines have been selected for editorial development — each "
        "centred on a named custodian, a specific place, and a practice that "
        "can still be witnessed."
    ),
    "angle": (
        "We focus on endangered knowledge systems and under-documented living "
        "practices surviving despite standardisation, policy change and "
        "generational drift. The genocide is context, never the lens; the state "
        "is never cast as villain. Every storyline centres a named custodian, a "
        "specific place, and a practice that can still be witnessed."
    ),
    "angle_why": (
        "This is the equitable, co-created Africa–Europe cultural work the "
        "Goethe programme exists to fund: a Rwandan cultural partner and a "
        "European publication building a long-form visual record of living "
        "Rwandan traditions — made with custodians, not about them, and "
        "presented in both Rwanda and Europe."
    ),
    "categories": [
        ("Food", 4),
        ("Myth & Religion", 3),
        ("Local communities", 2),
        ("Arts & Culture", 2),
        ("Geopolitics", 1),
    ],
}

STORIES = [
    Story(
        headline="What They Carry",
        subtitle="Batwa Forest Knowledge",
        hook=(
            "Batwa elders around Musanze hold centuries-deep ecological knowledge "
            "of the Nyungwe and Virunga forests — plant identification, animal "
            "behaviour, spatial memory — now sought by conservation biologists "
            "before it disappears. ~25,000–30,000 Batwa, officially ‘historically "
            "marginalised’."
        ),
        proposal=(
            "INLAND embeds with elders on a plant walk and records the forest "
            "songs — documenting the knowledge itself rather than the (already "
            "filmed) displacement. The narrative asks what becomes of forest "
            "knowledge when there is no forest left to use it in."
        ),
        category="Local communities",
        tags=["Endangered Knowledge", "Indigenous"],
        pull_quote="Not what was taken — what they still carry.",
        places="Musanze · Kibingo · forest margins",
        reference="AIMPO; INLAND Rwanda research corpus",
    ),
    Story(
        headline="The Rites That Won’t Die",
        subtitle="Kubandwa / Ryangombe",
        hook=(
            "A pre-colonial spirit-possession system centred on Ryangombe that "
            "missionaries tried to erase, surviving underground in a 92.8% "
            "Christian country. A 2021 academic recorded students attending rites "
            "‘even in the centre of Kigali’."
        ),
        proposal=(
            "INLAND follows the practice through discreet, consent-led fieldwork "
            "and academic introductions — documenting an adaptive ritual that "
            "crosses linguistic, political and temporal frontiers. If confirmed "
            "in the field, the edition’s signature story."
        ),
        category="Myth & Religion",
        tags=["Endangered Knowledge", "Sacred & Hidden"],
        pull_quote="Gone underground, not extinct.",
        places="Kigali · discreet rural sites",
        reference="Pennacini (Turin); Bizimana, RCSP Rwanda 2018",
    ),
    Story(
        headline="The Hospital Sends You to a Healer",
        subtitle="Uburozi and Traditional Medicine",
        hook=(
            "In Kinigi, Musanze — ‘the core location for healers’ — hospitals "
            "refer patients to traditional healers for uburozi, a sorcery/poisoning "
            "illness conventional medicine doesn’t recognise. ~5,000 healers "
            "nationally; a 2022 Traditional Medicine policy exists."
        ),
        proposal=(
            "INLAND shadows a referral from clinic to healer and back, with "
            "portraits of both — documenting a state-acknowledged parallel "
            "healthcare system operating in the open and almost entirely "
            "unreported internationally."
        ),
        category="Myth & Religion",
        tags=["State & Tradition", "Endangered Knowledge"],
        pull_quote="A parallel medicine, acknowledged and unseen.",
        places="Kinigi · Burera · Bushozi · Musenyi",
        reference="BMC Compl. Med. & Therapies 2021; Rwanda FDA 2022",
    ),
    Story(
        headline="Crushed by Foot",
        subtitle="Urwagwa, the Banana Beer of Rites",
        hook=(
            "Banana beer is ritual infrastructure — the medium of marriage "
            "negotiations, ancestral offerings and reconciliation. Named brewer "
            "Agatha Tuyisenge (Mushubati, Rutsiro), mother of five, built a "
            "livelihood from inherited brewing knowledge."
        ),
        proposal=(
            "INLAND follows one full brewing cycle toward a wedding — variety "
            "selection, foot-crushing, sorghum fermentation — and the ceremonial "
            "economy around the beer. A named protagonist with a full story arc."
        ),
        category="Food",
        tags=["Endangered Knowledge", "Ritual"],
        pull_quote="Not a drink — the medium ceremonies happen in.",
        places="Mushubati · Rutsiro District",
        reference="SENS Magazine (Tuyisenge profile)",
    ),
    Story(
        headline="The Last Sorghum Beer",
        subtitle="Ikigage on the Volcano Slopes",
        hook=(
            "A sorghum beer of feasts, dowries and harvests near Kinigi and "
            "Musanze, now on the Slow Food Ark of Taste — declining as sorghum "
            "is squeezed out by the Crop Intensification Programme."
        ),
        proposal=(
            "INLAND documents a household that still grows sorghum and brews "
            "ikigage for weddings — a compressed portrait of what agricultural "
            "policy does to cultural memory."
        ),
        category="Food",
        tags=["Endangered Knowledge", "State & Tradition"],
        pull_quote="A beer on the Ark of Taste, brewed for one more wedding.",
        places="Kinigi · Musanze",
        reference="Slow Food Ark of Taste; Slow Food Rwanda",
    ),
    Story(
        headline="A Harvest Without Its Secrets",
        subtitle="Umuganura and the Abiru",
        hook=(
            "Umuganura was a sacred rite of ubwiru guarded by the hereditary "
            "Abiru, needing sacred sorghum from Huro (Gakenke). Abolished by the "
            "Belgians in 1925; revived in 2011 as a secular holiday — without the "
            "Abiru, the plants, or the secret language."
        ),
        proposal=(
            "INLAND shoots the intact-but-vacant ritual site at Huro against the "
            "August state festival — a ceremony hollowed of meaning and "
            "repurposed for nation-building. A story about what revival leaves out."
        ),
        category="Myth & Religion",
        tags=["History", "Sacred & Hidden"],
        pull_quote="The site survived. The secret did not.",
        places="Huro · Bumbogo · Gakenke District",
        reference="ubwiru / Abiru historiography",
    ),
    Story(
        headline="The Morality of Milk",
        subtitle="Kivuguto and the Inkongoro Gourd",
        hook=(
            "Fermented milk made in handmade inkongoro gourds carrying unique "
            "indigenous microbial cultures — the moral order encoded in how milk "
            "is stored, fermented and shared, now replaced by steel tanks and "
            "controlled starters."
        ),
        proposal=(
            "INLAND builds a diptych: morning milking and gourd fermentation in a "
            "pastoral household, against a modern dairy plant making packaged "
            "‘yogurt drinks’. Elders speak to the symbolism of the gourd."
        ),
        category="Food",
        tags=["Endangered Knowledge", "Ritual"],
        pull_quote="A microbial culture is also a moral one.",
        places="Pastoral households · Musanze / Gishwati",
        reference="INLAND Rwanda research corpus",
    ),
    Story(
        headline="Life After the Forest",
        subtitle="The Gishwati Resettlement",
        hook=(
            "Families moved off the Gishwati–Mukura forests now farm ~0.5 ha on "
            "steep, flood-prone slopes in Nyabihu, Rutsiro and Rubavu. Some, like "
            "Laurent Hategekimana, patrol as unpaid auxiliaries the forest margins "
            "their families were removed from."
        ),
        proposal=(
            "INLAND walks the forest margin with a patroller and spends time in "
            "the resettlement villages — the long, quiet afterlife of displacement "
            "on a landscape remade by ranching, floods and restoration."
        ),
        category="Local communities",
        tags=["Land & Displacement", "Nature"],
        pull_quote="Guarding the forest that displaced them.",
        places="Nyabihu · Rutsiro · Rubavu · Bigogwe",
        reference="Rainforest Journalism Fund; Pulitzer Center",
    ),
    Story(
        headline="Thirty Years On",
        subtitle="The Reconciliation Coffee Cooperatives",
        hook=(
            "Cooperatives built post-1994 to put Hutu and Tutsi farmers in the "
            "same washing stations — Maraba (Huye) and Dukunde Kawa Musasa "
            "(Gakenke). The origin story is told; the present is unfilmed."
        ),
        proposal=(
            "INLAND follows one or two families through a harvest — picking, "
            "sorting, washing — asking quietly what changed in a generation and "
            "what didn’t. The founders are middle-aged now; their children work "
            "side by side."
        ),
        category="Food",
        tags=["History", "Reconciliation"],
        pull_quote="The origin story is told. The present is not.",
        places="Maraba (Huye) · Dukunde Kawa (Gakenke)",
        reference="Alliance for Coffee Excellence; cooperative profiles",
    ),
    Story(
        headline="The Last Maestro",
        subtitle="Vianey Mushabizi and the Inanga",
        hook=(
            "A single instrument and the man who may be the end of its line: "
            "Vianey Mushabizi (Gatagara), possibly the last living master of the "
            "inanga zither — the sound of a pre-radio Rwanda held in one pair of "
            "hands."
        ),
        proposal=(
            "INLAND records his repertoire and tuning, and asks the question of "
            "transmission: is there an apprentice, or does the line end here? A "
            "portrait in sound and image."
        ),
        category="Arts & Culture",
        tags=["Endangered Knowledge", "Music"],
        pull_quote="A whole tradition in one pair of hands.",
        places="Gatagara",
        reference="INLAND Rwanda research corpus",
    ),
    Story(
        headline="Geometry of the Widows",
        subtitle="Imigongo and the Market",
        hook=(
            "Geometric cow-dung relief paintings made by genocide widows in "
            "Kakira, Eastern Province — now commercialised. A tension between "
            "cultural preservation, livelihood and market co-option."
        ),
        proposal=(
            "INLAND documents the making and the market in one place — who paints, "
            "who sells, who profits when a domestic craft scales into a souvenir."
        ),
        category="Arts & Culture",
        tags=["Craft", "History"],
        pull_quote="Who profits when a tradition scales?",
        places="Kakira · Eastern Province",
        reference="INLAND Rwanda research corpus",
    ),
    Story(
        headline="The Data Centre and the Hill",
        subtitle="Two Relationships to Electricity",
        hook=(
            "A windowless server warehouse runs 24 hours at full load a few "
            "kilometres from villages where electricity arrived three years ago. "
            "The guard at the gate is from such a village."
        ),
        proposal=(
            "INLAND anchors the development showcase in one human — the guard — "
            "to portray two relationships to electricity, land and the future "
            "coexisting within a few kilometres, without abstraction."
        ),
        category="Geopolitics",
        tags=["Modernisation", "Land"],
        pull_quote="Two centuries, a few kilometres apart.",
        places="Eastern Province / peri-Kigali",
        reference="INLAND Rwanda research corpus",
    ),
]

TRACKS = [
    {
        "name": "The Musanze Cluster",
        "color": (0x2C, 0x5F, 0x2D),
        "stories": "Stories 01 + 03 + 05",
        "subtitle": "Knowledge, healing & ikigage",
        "desc": (
            "A single northern base yields three nested stories: Batwa "
            "forest knowledge, the hospital-to-healer uburozi referral, and "
            "the last sorghum beer. One deployment, three custodian-led "
            "narratives around Musanze and Kinigi."
        ),
        "geography": "Musanze → Kinigi → forest margins",
        "season": "Flexible (dry season for access)",
        "story_count": "3 stories, 1 field base",
    },
    {
        "name": "Brewing & Land",
        "color": (0xB8, 0x50, 0x42),
        "stories": "Stories 04 + 08 (+ 07)",
        "subtitle": "Ritual food & displacement",
        "desc": (
            "A western cluster binding the urwagwa brewing cycle in Rutsiro, "
            "the Gishwati resettlement in Nyabihu, and the kivuguto milk "
            "diptych. Ritual food and the afterlife of displacement on the "
            "same hillsides."
        ),
        "geography": "Rutsiro → Nyabihu → Rubavu",
        "season": "Toward a wedding / harvest",
        "story_count": "2–3 stories, 1 deployment",
    },
    {
        "name": "Coffee & Memory",
        "color": (0x5B, 0x2C, 0x6F),
        "stories": "Stories 09 + 06 (+ 11)",
        "subtitle": "Cooperatives, harvest & craft",
        "desc": (
            "A southern/eastern cluster: the reconciliation coffee "
            "cooperatives thirty years on, the hollowed Umuganura harvest at "
            "Huro, and the imigongo widows. Memory, harvest and craft across "
            "Huye, Gakenke and the East."
        ),
        "geography": "Huye → Gakenke → Eastern Province",
        "season": "Coffee harvest; August (Umuganura)",
        "story_count": "2–3 stories, seasonal windows",
    },
]


def build_pptx(output_path: str = "INLAND_Rwanda_Pitch.pptx"):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    import os

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG = RGBColor(0xF9, 0xF4, 0xEE)
    FG = RGBColor(0x00, 0x00, 0x00)
    ACCENT = RGBColor(0x40, 0x05, 0x05)
    MUTED = RGBColor(0x5C, 0x5C, 0x5C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x0E, 0x0E, 0x0E)
    MUTED_WHITE = RGBColor(0x99, 0x99, 0x99)
    REF_COLOR = RGBColor(0x85, 0x7F, 0x76)

    HEADING = "Baskerville"
    BODY = "Baskerville"
    LABEL = "Fraunces"

    script_dir = os.path.dirname(os.path.abspath(__file__))
    logo_black = os.path.join(script_dir, "inland-logo-black.png")
    logo_white = os.path.join(script_dir, "inland-logo-white.png")

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def text(slide, left, top, w, h, txt, font, size, color,
             bold=False, align=PP_ALIGN.LEFT, italic=False, spacing=None):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = txt
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = align
        p.line_spacing = Pt(spacing) if spacing else Pt(int(size * 1.6))
        return box

    def line(slide, left, top, w, color=FG):
        shape = slide.shapes.add_connector(1, Inches(left), Inches(top), Inches(left + w), Inches(top))
        shape.line.color.rgb = color
        shape.line.width = Pt(0.75)

    def vline(slide, left, top, h, color=ACCENT):
        shape = slide.shapes.add_connector(1, Inches(left), Inches(top), Inches(left), Inches(top + h))
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

    def add_tag_pill(slide, left, top, tag_text, font=LABEL, size=9, color=MUTED):
        w = len(tag_text) * 0.075 + 0.28
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(0.3))
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        spPr = box._element.find(qn('a:spPr'))
        if spPr is None:
            spPr = box._element.makeelement(qn('a:spPr'), {})
            box._element.append(spPr)
        ln = spPr.makeelement(qn('a:ln'), {'w': '6350'})
        solidFill = ln.makeelement(qn('a:solidFill'), {})
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': 'B0A89E'})
        solidFill.append(srgb)
        ln.append(solidFill)
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'solid'}))
        spPr.append(ln)
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
        else:
            prstGeom = spPr.makeelement(qn('a:prstGeom'), {'prst': 'roundRect'})
            prstGeom.append(prstGeom.makeelement(qn('a:avLst'), {}))
            spPr.insert(0, prstGeom)
        return w

    # Title
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    logo_w, logo_h = Inches(3.5), Inches(0.735)
    logo_left = (prs.slide_width - logo_w) // 2
    if os.path.exists(logo_black):
        sl.shapes.add_picture(logo_black, logo_left, Inches(1.8), logo_w, logo_h)
    else:
        text(sl, 1, 1.6, 11.3, 0.8, "INLAND", LABEL, 48, FG, bold=True, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.9, 11.3, 1.4, "Rwanda", HEADING, 72, FG, align=PP_ALIGN.CENTER)
    line(sl, 5.5, 4.6, 2.3)
    text(sl, 1, 5.0, 11.3, 0.6, "An Africa–Europe Cultural Co-Production", BODY, 20, MUTED, align=PP_ALIGN.CENTER)
    text(sl, 1, 6.2, 11.3, 0.4, "12 Selected Stories  ·  For the Goethe-Institut Africa–Europe Partnerships for Culture  ·  2026",
         LABEL, 11, MUTED, align=PP_ALIGN.CENTER)

    # Contents
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "CONTENTS", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Overview", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    for j, (num, title) in enumerate([
        ("01", "About INLAND"), ("02", "Rwanda — What Survived"),
        ("03", "Our Angle & Categories"), ("04", "The 12 Storylines"),
        ("05", "Production Tracks"), ("06", "The Partnership")]):
        y = 2.85 + j * 0.62
        text(sl, 0.8, y, 1.0, 0.5, num, HEADING, 24, ACCENT)
        text(sl, 2.0, y + 0.03, 8, 0.5, title, BODY, 20, FG)
        if j < 5:
            line(sl, 0.8, y + 0.52, 11.5, RGBColor(0xD8, 0xD2, 0xCA))

    # About
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "ABOUT", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "About INLAND", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.0, INTRO["about"], BODY, 17, FG, spacing=27)
    text(sl, 0.8, 5.3, 7.5, 1.5, INTRO["about_method"], BODY, 14, MUTED)
    text(sl, 9.5, 2.9, 3.2, 1.5, "Intimacy with\ncustodians, not\nthematic surveys.", HEADING, 24, ACCENT, italic=True)

    # Research
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THE RESEARCH", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Rwanda — What Survived", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 2.0, INTRO["research"], BODY, 17, FG, spacing=27)
    text(sl, 0.8, 5.5, 7.5, 1.5, INTRO["research_scope"], BODY, 14, MUTED)
    text(sl, 9.7, 2.7, 3, 1.2, "20", HEADING, 80, ACCENT)
    text(sl, 9.7, 3.8, 3, 0.4, "story leads researched", LABEL, 11, MUTED)
    text(sl, 9.7, 4.5, 3, 1.2, "12", HEADING, 80, ACCENT)
    text(sl, 9.7, 5.6, 3, 0.4, "selected for this round", LABEL, 11, MUTED)

    # Angle & Categories
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "EDITORIAL ANGLE", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Our Angle", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    text(sl, 0.8, 2.9, 7.5, 1.8, INTRO["angle"], BODY, 17, FG, spacing=27)
    text(sl, 0.8, 5.3, 7.5, 1.5, INTRO["angle_why"], BODY, 14, MUTED)
    text(sl, 9.7, 2.7, 3.5, 0.4, "CATEGORIES", LABEL, 11, MUTED)
    line(sl, 9.7, 3.15, 2.9)
    for j, (cat, count) in enumerate(INTRO["categories"]):
        y = 3.35 + j * 0.55
        text(sl, 9.7, y, 2.5, 0.4, cat, BODY, 16, FG)
        text(sl, 12.1, y, 1, 0.4, str(count), BODY, 16, ACCENT, align=PP_ALIGN.RIGHT, bold=True)
        if j < len(INTRO["categories"]) - 1:
            line(sl, 9.7, y + 0.45, 2.9, RGBColor(0xD8, 0xD2, 0xCA))

    # Breaker
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.0, 11.3, 0.4, "PART TWO", LABEL, 11, MUTED_WHITE, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.8, 11.3, 1.2, "The Storylines", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    text(sl, 2.5, 4.6, 8.3, 1.2,
         "Twelve selected stories across communities, ritual, food, ecology, "
         "craft and music. Each centres a named custodian, a specific place, "
         "and a practice that can still be witnessed.", BODY, 18, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # Story slides
    PILL = RGBColor(0x70, 0x70, 0x70)
    for i, s in enumerate(STORIES):
        sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
        text(sl, 11.0, 0.4, 2, 1.0, f"{i+1:02d}", HEADING, 64, ACCENT, align=PP_ALIGN.RIGHT)
        text(sl, 0.8, 0.6, 5, 0.4, s.category.upper(), LABEL, 11, MUTED)
        text(sl, 0.8, 1.6, 10, 1.2, s.headline, HEADING, 42, FG)
        text(sl, 0.8, 2.8, 10, 0.8, s.subtitle, HEADING, 26, MUTED, italic=True)
        line(sl, 0.8, 3.7, 11.5)
        text(sl, 0.8, 3.9, 2, 0.25, "THE STORY", LABEL, 8, PILL)
        text(sl, 0.8, 4.15, 5.8, 1.8, s.hook, BODY, 14, FG, spacing=22)
        vline(sl, 7.0, 3.9, 2.4, RGBColor(0xD8, 0xD2, 0xCA))
        text(sl, 7.4, 3.9, 2, 0.25, "OUR PROPOSAL", LABEL, 8, PILL)
        text(sl, 7.4, 4.15, 5.2, 1.8, s.proposal, BODY, 14, FG, spacing=22)
        text(sl, 0.8, 6.55, 7.2, 0.35, f"Places: {s.places}", LABEL, 10, MUTED)
        text(sl, 0.8, 6.92, 7.2, 0.3, f"Sources: {s.reference}", LABEL, 9, REF_COLOR)
        tag_x = 12.5
        for tag in reversed(s.tags):
            w = add_tag_pill(sl, 0, 6.85, tag)
            tag_x -= (w + 0.15)
            sl.shapes[-1].left = Inches(tag_x + 0.15)

    # Tracks breaker
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.0, 11.3, 0.4, "OUR SUGGESTION", LABEL, 11, MUTED_WHITE, align=PP_ALIGN.CENTER)
    text(sl, 1, 2.8, 11.3, 1.2, "Production Tracks", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    text(sl, 2.5, 4.6, 8.3, 1.2,
         "From the twelve storylines we propose three production tracks for a "
         "pilot — each grouping stories that share geography and field timing "
         "into one deployable trip with the Rwandan partner.", BODY, 18, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # Tracks overview
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THREE TRACKS", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 8, 1.0, "Production Tracks", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    cw, gap, sx, top = 3.7, 0.35, 0.8, 2.9
    for j, track in enumerate(TRACKS):
        x = sx + j * (cw + gap)
        tc = RGBColor(*track["color"])
        bar = sl.shapes.add_connector(1, Inches(x), Inches(top), Inches(x + cw), Inches(top))
        bar.line.color.rgb = tc; bar.line.width = Pt(3)
        text(sl, x, top + 0.15, cw, 0.5, track["name"], HEADING, 23, tc)
        text(sl, x, top + 0.62, cw, 0.3, track["stories"], LABEL, 9, MUTED)
        text(sl, x, top + 0.95, cw, 1.7, track["desc"], BODY, 13, FG, spacing=18)
        meta = f"Geography: {track['geography']}\nSeason: {track['season']}\n{track['story_count']}"
        text(sl, x, top + 2.9, cw, 1.0, meta, LABEL, 8, REF_COLOR)

    # Partnership slide
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG)
    text(sl, 0.8, 0.5, 4, 0.4, "THE PARTNERSHIP", LABEL, 11, MUTED)
    text(sl, 0.8, 1.3, 11, 1.0, "An Equitable Co-Production", HEADING, 40, FG)
    line(sl, 0.8, 2.5, 11.5)
    pairs = [
        ("European partner", "INLAND — yearly publication (Italy / EU). Editorial, photography and publishing; presentation of the work in Europe."),
        ("African partner", "A Rwandan cultural organisation (lead candidate: Red Rocks Initiatives, Musanze) — co-creation, community access, custodian liaison, presentation in Rwanda."),
        ("Funding path", "Mobility grant (≤€4,000) funds the scouting trip; the Visual Arts grant funds collaborative production and a public presentation in Rwanda and Europe."),
        ("The principle", "Made with custodians, not about them — equitable Africa–Europe co-creation, exactly the dialogue the Goethe programme exists to fund."),
    ]
    for j, (h, b) in enumerate(pairs):
        y = 2.9 + j * 1.05
        text(sl, 0.8, y, 3.0, 0.4, h, LABEL, 12, ACCENT, bold=True)
        text(sl, 4.0, y - 0.02, 8.6, 1.0, b, BODY, 14, FG, spacing=19)

    # Closing
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, DARK)
    text(sl, 1, 2.2, 11.3, 1.5, "12 Storylines.\nOne Rwanda.", HEADING, 52, WHITE, align=PP_ALIGN.CENTER)
    line(sl, 5.8, 4.2, 1.7, WHITE)
    lw, lh = Inches(2), Inches(0.42)
    if os.path.exists(logo_white):
        sl.shapes.add_picture(logo_white, (prs.slide_width - lw) // 2, Inches(5.0), lw, lh)
        text(sl, 1, 5.6, 11.3, 0.4, "An Africa–Europe co-production proposal", LABEL, 14, MUTED_WHITE, align=PP_ALIGN.CENTER)
    else:
        text(sl, 1, 4.9, 11.3, 0.6, "An Africa–Europe co-production proposal by INLAND", LABEL, 14, MUTED_WHITE, align=PP_ALIGN.CENTER)

    # strip shadows
    for sl in prs.slides:
        for shape in sl.shapes:
            spPr = shape._element.find(qn('p:spPr'))
            if spPr is not None:
                for ef in spPr.findall(qn('a:effectLst')):
                    spPr.remove(ef)
                spPr.append(spPr.makeelement(qn('a:effectLst'), {}))

    full = os.path.join(script_dir, output_path)
    prs.save(full)
    print(f"Saved PPTX: {full} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build_pptx()
